#!/usr/bin/env python3
"""Continue BagDrop PPT V4 Final - Slides 15-20 (FINAL)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Load existing presentation
prs = Presentation(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4_Final.pptx"
)

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 215, 0)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_GRAY = RGBColor(120, 120, 120)
GREEN = RGBColor(52, 168, 83)
RED = RGBColor(234, 67, 53)

HEADING_FONT = "Arial Black"
BODY_FONT = "Montserrat"


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_heading(
    slide, left, top, width, height, text, font_size, color, align=PP_ALIGN.LEFT
):
    return add_text(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size,
        color,
        bold=True,
        font=HEADING_FONT,
        align=align,
    )


def set_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================
# SLIDE 15: RISK MITIGATION
# ============================================
s15 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s15, WHITE)
add_heading(
    s15, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), "RISK MITIGATION", 56, BLACK
)

risks = [
    ("LOW BOOKINGS", "Aggressive referral program\nDynamic pricing adjustments"),
    ("PARTNER DEFAULTS", "Weekly settlements\n₹5K security deposit"),
    ("FRAUD / THEFT", "Insurance up to ₹25K/bag\nPartner verification & CCTV"),
    ("COMPETITION", "Focus on USP (delivery)\nPartner lock-in agreements"),
    ("PAYMENT DISPUTES", "Clear T&C\nEscrow for high-value items"),
]

for i, (risk, mitigation) in enumerate(risks):
    y = Inches(1.4 + i * 0.95)

    # Risk
    add_rect(s15, Inches(0.5), y, Inches(0.15), Inches(0.6), RED)
    add_text(s15, Inches(0.8), y, Inches(3.5), Inches(0.4), risk, 16, RED, bold=True)

    # Arrow
    add_text(s15, Inches(4.5), y, Inches(0.5), Inches(0.4), "→", 20, YELLOW, bold=True)

    # Mitigation
    add_text(s15, Inches(5.2), y, Inches(7.5), Inches(0.6), mitigation, 13, DARK_GRAY)

# Insurance banner
add_rect(s15, Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.8), BLACK)
insurance_txt = "INSURANCE: Public Liability ₹10L | Storage Insurance ₹25K/bag | Employee Health + Accident"
add_text(
    s15,
    Inches(1),
    Inches(6.35),
    Inches(11.333),
    Inches(0.5),
    insurance_txt,
    13,
    YELLOW,
    align=PP_ALIGN.CENTER,
)

add_notes(
    s15,
    "We've identified risks and mitigation strategies. Low bookings? Referral programs. Partner defaults? Weekly settlements with deposits. Theft? Insurance and CCTV. Competition? Our delivery USP. Payment disputes? Clear terms and escrow.",
)

# ============================================
# SLIDE 16: EXIT STRATEGY
# ============================================
s16 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s16, BLACK)
add_heading(
    s16, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "EXIT STRATEGY", 64, YELLOW
)

exits = [
    ("STRATEGIC ACQUISITION", "70% probability", "MakeMyTrip, Goibibo, Delhivery"),
    ("MERGER (LOGISTICS)", "15% probability", "Delhivery, Dunzo, Porter"),
    ("TECHNOLOGY LICENSING", "10% probability", "White-label platform sales"),
    ("IPO (LONG-TERM)", "5% probability", "Year 7+ horizon"),
]

for i, (exit_type, prob, example) in enumerate(exits):
    y = Inches(1.5 + i * 1.1)

    add_rect(s16, Inches(0.5), y, Inches(0.15), Inches(0.8), YELLOW)
    add_text(
        s16, Inches(0.8), y, Inches(4), Inches(0.4), exit_type, 18, WHITE, bold=True
    )
    add_text(s16, Inches(5), y, Inches(2.5), Inches(0.4), prob, 14, GREEN)
    add_text(s16, Inches(7.8), y, Inches(5), Inches(0.4), example, 13, LIGHT_GRAY)

# Valuation
add_rect(s16, Inches(0.5), Inches(6), Inches(6), Inches(1), YELLOW)
add_text(
    s16,
    Inches(0.7),
    Inches(6.15),
    Inches(5.6),
    Inches(0.5),
    "TARGET VALUATION",
    16,
    BLACK,
    bold=True,
)
add_text(
    s16,
    Inches(0.7),
    Inches(6.6),
    Inches(5.6),
    Inches(0.35),
    "₹50+ Crore by Year 5",
    22,
    BLACK,
    bold=True,
)

add_notes(
    s16,
    "4 exit paths. Most likely: acquisition by MakeMyTrip/Goibibo/Delhivery (70%). Merger with logistics (15%). Tech licensing (10%). IPO long-term (5%). Target: ₹50+ Cr valuation by Year 5.",
)

# ============================================
# SLIDE 17: MILESTONES / ROADMAP
# ============================================
s17 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s17, YELLOW)
add_heading(
    s17, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "ROADMAP", 64, BLACK
)

# Timeline
milestones = [
    ("Q1 2026", "App launch\n50 partners in Gujarat"),
    ("Q2 2026", "5,000 monthly bookings\nExpand to 3 cities"),
    ("Q3 2026", "10,000 bookings\nDelivery service live"),
    ("Q4 2026", "30,000 bookings\nBreak-even achieved"),
    ("YEAR 2", "50 cities\n500 partners\n₹1.27Cr revenue"),
    ("YEAR 3", "150 cities\n1,500 partners\n₹3.82Cr revenue"),
]

for i, (time, milestone) in enumerate(milestones):
    x = Inches(0.5 + i * 2.1)
    color = BLACK if i % 2 == 0 else WHITE
    txt_color = YELLOW if i % 2 == 0 else BLACK

    add_rect(s17, x, Inches(1.8), Inches(1.9), Inches(1.8), color)
    add_text(
        s17,
        x,
        Inches(2),
        Inches(1.9),
        Inches(0.6),
        time,
        18,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s17,
        x,
        Inches(2.7),
        Inches(1.9),
        Inches(0.8),
        milestone,
        11,
        WHITE if i % 2 == 0 else DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )

# Vision
add_rect(s17, Inches(0.5), Inches(4), Inches(12.333), Inches(1), BLACK)
add_text(
    s17,
    Inches(1),
    Inches(4.2),
    Inches(11.333),
    Inches(0.6),
    "VISION: India's #1 luggage storage network with smart locker integrations at all major transit hubs",
    18,
    WHITE,
    align=PP_ALIGN.CENTER,
)

# Year 5 target
add_heading(
    s17,
    Inches(0.8),
    Inches(5.5),
    Inches(11),
    Inches(0.8),
    "YEAR 5: EXIT READINESS",
    32,
    BLACK,
)
targets = [("₹50+ Cr", "Valuation"), ("100+", "Cities"), ("2,000+", "Partners")]
for i, (val, label) in enumerate(targets):
    x = Inches(0.8 + i * 4)
    add_text(s17, x, Inches(6.2), Inches(3.5), Inches(0.5), val, 28, BLACK, bold=True)
    add_text(s17, x, Inches(6.7), Inches(3.5), Inches(0.3), label, 12, DARK_GRAY)

add_notes(
    s17,
    "5-year roadmap. Q1: App launch. Q4: Break-even at 30K bookings. Year 2: 50 cities, ₹1.27Cr. Year 3: 150 cities, ₹3.82Cr. Year 5: Exit readiness with ₹50+ Cr valuation.",
)

# ============================================
# SLIDE 18: GOVERNMENT SUPPORT
# ============================================
s18 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s18, WHITE)
add_heading(
    s18,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.8),
    "GOVERNMENT SUPPORT",
    56,
    BLACK,
)

schemes = [
    ("GUJARAT STARTUP POLICY", "Up to ₹30 Lakh seed funding"),
    ("STARTUP INDIA", "3-year tax holiday"),
    ("GST BENEFITS", "Input credit on tech expenses"),
    ("MUDRA / MSME LOANS", "₹10-50 Lakh low-interest debt"),
    ("iCREATE / AIC", "Free mentoring & co-working"),
]

for i, (scheme, benefit) in enumerate(schemes):
    y = Inches(1.4 + i * 0.9)

    add_rect(s18, Inches(0.5), y, Inches(0.15), Inches(0.6), GREEN)
    add_text(
        s18, Inches(0.8), y, Inches(4.5), Inches(0.4), scheme, 16, DARK_GRAY, bold=True
    )
    add_text(
        s18,
        Inches(0.8),
        y + Inches(0.4),
        Inches(6),
        Inches(0.4),
        benefit,
        13,
        LIGHT_GRAY,
    )

# Kutch advantage
add_rect(s18, Inches(7), Inches(2), Inches(5.8), Inches(2.5), YELLOW)
add_text(
    s18,
    Inches(7.2),
    Inches(2.3),
    Inches(5.4),
    Inches(0.5),
    "KUTCH ADVANTAGE",
    20,
    BLACK,
    bold=True,
)
add_text(
    s18,
    Inches(7.2),
    Inches(2.9),
    Inches(5.4),
    Inches(1.5),
    "Gujarat government actively supports tourism & logistics startups\n\nAccess to iCreate, AIC mentoring\n\nGujarat Startup Policy benefits",
    14,
    DARK_GRAY,
)

add_notes(
    s18,
    "Leveraging government support: Gujarat Startup Policy (₹30L), Startup India tax holiday, GST benefits, MUDRA loans, iCreate/AIC mentoring. Kutch location gives us access to active tourism startup support.",
)

# ============================================
# SLIDE 19: CALL TO ACTION
# ============================================
s19 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s19, BLACK)
add_heading(
    s19, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "JOIN US", 64, YELLOW
)

# Investment box
add_rect(s19, Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.5), YELLOW)
add_text(
    s19,
    Inches(0.7),
    Inches(1.7),
    Inches(5.1),
    Inches(0.5),
    "SEED FUNDING",
    18,
    BLACK,
    bold=True,
)
add_text(
    s19,
    Inches(0.7),
    Inches(2.3),
    Inches(5.1),
    Inches(1),
    "₹36",
    72,
    BLACK,
    bold=True,
    font=HEADING_FONT,
)
add_text(
    s19,
    Inches(0.7),
    Inches(3.3),
    Inches(5.1),
    Inches(0.4),
    "LAKH",
    24,
    BLACK,
    bold=True,
)
add_text(
    s19,
    Inches(0.7),
    Inches(3.8),
    Inches(5.1),
    Inches(0.3),
    "(~$430K USD)",
    12,
    DARK_GRAY,
)

# Why invest
why_points = [
    "✓ Solves real problem (70% travelers)",
    "✓ Unique USP: Forgotten luggage delivery",
    "✓ Asset-light, scalable model",
    "✓ Experienced local team",
    "✓ Government support & benefits",
    "✓ Clear path to profitability",
]

for i, point in enumerate(why_points):
    y = Inches(1.5 + i * 0.6)
    add_text(s19, Inches(6.5), y, Inches(6.5), Inches(0.5), point, 14, WHITE)

# ROI
add_rect(s19, Inches(0.5), Inches(4.5), Inches(5.5), Inches(1.2), RGBColor(40, 40, 40))
add_text(
    s19,
    Inches(0.7),
    Inches(4.7),
    Inches(5.1),
    Inches(0.4),
    "EXPECTED RETURNS",
    14,
    YELLOW,
    bold=True,
)
add_text(
    s19,
    Inches(0.7),
    Inches(5.2),
    Inches(5.1),
    Inches(0.4),
    "5-10x for seed investors",
    20,
    WHITE,
    bold=True,
)

# Contact
add_rect(s19, Inches(0.5), Inches(6), Inches(12.333), Inches(1.2), YELLOW)
add_text(
    s19,
    Inches(1),
    Inches(6.2),
    Inches(11.333),
    Inches(0.6),
    "TeamBack | Kutch, Gujarat | January 2026",
    20,
    BLACK,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_text(
    s19,
    Inches(1),
    Inches(6.9),
    Inches(11.333),
    Inches(0.3),
    "Ready to transform travel in India",
    14,
    DARK_GRAY,
    align=PP_ALIGN.CENTER,
)

add_notes(
    s19,
    "Call to action: ₹36 Lakh seed funding. Expected 5-10x returns. Why invest? Real problem, unique USP, scalable model, experienced team, government support, clear exit path. Join TeamBack from Kutch, Gujarat!",
)

# ============================================
# SLIDE 20: THANK YOU / Q&A
# ============================================
s20 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s20, BLACK)

add_heading(
    s20,
    Inches(0),
    Inches(2),
    Inches(13.333),
    Inches(1.5),
    "THANK YOU",
    96,
    YELLOW,
    align=PP_ALIGN.CENTER,
)
add_text(
    s20,
    Inches(0),
    Inches(3.8),
    Inches(13.333),
    Inches(0.8),
    "QUESTIONS & DISCUSSION",
    32,
    WHITE,
    align=PP_ALIGN.CENTER,
)
add_text(
    s20,
    Inches(0),
    Inches(5),
    Inches(13.333),
    Inches(1),
    "BagDrop: India's Smart Luggage Storage Network\nTeamBack | Kutch, Gujarat",
    20,
    LIGHT_GRAY,
    align=PP_ALIGN.CENTER,
)

# Yellow accent
add_rect(s20, Inches(4), Inches(6.5), Inches(5.333), Inches(0.15), YELLOW)

add_notes(
    s20,
    "Thank you for your attention! Happy to answer questions about BagDrop, our business model, financials, or partnership opportunities. Contact TeamBack in Kutch, Gujarat. Let's revolutionize travel in India!",
)

# SAVE FINAL
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4_Final.pptx"
)
print("=" * 60)
print("COMPLETE! All 20 slides created successfully!")
print("=" * 60)
print("Design: Swiss Style (Yellow/Black/White)")
print("Fonts: Arial Black + Montserrat")
print("File: BagDrop_Presentation_V4_Final.pptx")
print("=" * 60)
