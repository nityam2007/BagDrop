#!/usr/bin/env python3
"""
BagDrop Presentation V5 - IMPROVED VERSION
Fixes: Larger fonts on tables, better readability, restructured slides
Swiss Style: Yellow/Black/White
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 215, 0)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_GRAY = RGBColor(120, 120, 120)
GREEN = RGBColor(52, 168, 83)
RED = RGBColor(234, 67, 53)

HEADING_FONT = "Arial Black"
BODY_FONT = "Montserrat"


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_heading(
    slide, left, top, width, height, text, font_size, color, align=PP_ALIGN.LEFT
):
    return add_text(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size,
        color,
        bold=True,
        font=HEADING_FONT,
        align=align,
    )


def set_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================
# SLIDE 1: TITLE
# ============================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, WHITE)
add_rect(s1, Inches(0), Inches(2), Inches(2), Inches(3), YELLOW)
add_heading(s1, Inches(1.5), Inches(2), Inches(10), Inches(1.5), "BAGDROP", 96, BLACK)
add_text(
    s1,
    Inches(1.5),
    Inches(3.5),
    Inches(10),
    Inches(0.8),
    "SMART LUGGAGE STORAGE & DELIVERY",
    32,
    DARK_GRAY,
    bold=True,
)
add_text(
    s1,
    Inches(1.5),
    Inches(4.5),
    Inches(10),
    Inches(0.5),
    "Secure. Contactless. Affordable.",
    24,
    DARK_GRAY,
)
add_text(
    s1,
    Inches(1.5),
    Inches(6),
    Inches(10),
    Inches(0.5),
    "TeamBack | Kutch, Gujarat | January 2026",
    16,
    LIGHT_GRAY,
)
add_notes(
    s1,
    "Welcome! We're TeamBack from Kutch, Gujarat. BagDrop is India's smart luggage storage network solving a problem faced by 70% of travelers.",
)

# ============================================
# SLIDE 2: THE PROBLEM
# ============================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, WHITE)
add_rect(s2, Inches(10.5), Inches(0), Inches(2.833), Inches(1.8), YELLOW)
add_heading(
    s2, Inches(0.8), Inches(0.8), Inches(9), Inches(1), "THE PROBLEM", 64, BLACK
)
add_text(
    s2,
    Inches(0.8),
    Inches(1.8),
    Inches(9),
    Inches(0.6),
    "Travelers face a daily struggle with luggage",
    22,
    LIGHT_GRAY,
)

add_rect(
    s2, Inches(0.5), Inches(2.8), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245)
)
add_text(
    s2,
    Inches(0.8),
    Inches(3.1),
    Inches(5.4),
    Inches(1),
    "EARLY CHECKOUT, LATE DEPARTURE\nHotels force checkout but flights are hours away",
    18,
    DARK_GRAY,
)

add_rect(
    s2, Inches(6.8), Inches(2.8), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245)
)
add_text(
    s2,
    Inches(7.1),
    Inches(3.1),
    Inches(5.4),
    Inches(1),
    "HEAVY BAGS, RUINED PLANS\nDragging luggage kills the travel experience",
    18,
    DARK_GRAY,
)

add_rect(
    s2, Inches(0.5), Inches(4.5), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245)
)
add_text(
    s2,
    Inches(0.8),
    Inches(4.8),
    Inches(5.4),
    Inches(1),
    "CLOAKROOMS? OVERPRICED\nRailway lockers cost ₹100+/day, often full",
    18,
    DARK_GRAY,
)

add_rect(
    s2, Inches(6.8), Inches(4.5), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245)
)
add_text(
    s2,
    Inches(7.1),
    Inches(4.8),
    Inches(5.4),
    Inches(1),
    "FORGOT YOUR BAG?\nNo easy way to retrieve forgotten luggage",
    18,
    DARK_GRAY,
)

add_rect(s2, Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.6), BLACK)
add_text(
    s2,
    Inches(0.8),
    Inches(6.3),
    Inches(11.8),
    Inches(0.5),
    "70% of travelers have faced luggage hassles during their trip",
    20,
    YELLOW,
    bold=True,
)
add_notes(
    s2,
    "70% of travelers face luggage problems. Early checkout with late departure means hours of frustration.",
)

# ============================================
# SLIDE 3: THE SOLUTION
# ============================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, BLACK)
add_rect(s3, Inches(0), Inches(0), Inches(0.3), Inches(7.5), YELLOW)
add_heading(
    s3, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "THE SOLUTION", 64, YELLOW
)
add_text(
    s3,
    Inches(0.8),
    Inches(1.6),
    Inches(11),
    Inches(1),
    "Store your bags anywhere. Pick them up anytime.",
    32,
    WHITE,
    bold=True,
)

add_rect(s3, Inches(0.8), Inches(3), Inches(0.15), Inches(0.8), YELLOW)
add_text(
    s3,
    Inches(1.2),
    Inches(3),
    Inches(10),
    Inches(0.9),
    "NEIGHBORHOOD NETWORK\nWe partner with verified local shops, cafes & hotels",
    20,
    WHITE,
)

add_rect(s3, Inches(0.8), Inches(4.2), Inches(0.15), Inches(0.8), YELLOW)
add_text(
    s3,
    Inches(1.2),
    Inches(4.2),
    Inches(10),
    Inches(0.9),
    "INSTANT ACCESS\nBook via app, drop bag, get secure OTP or NFC token",
    20,
    WHITE,
)

add_rect(s3, Inches(0.8), Inches(5.4), Inches(0.15), Inches(0.8), YELLOW)
add_text(
    s3,
    Inches(1.2),
    Inches(5.4),
    Inches(10),
    Inches(0.9),
    "FORGOTTEN LUGGAGE? DELIVERED\nWe'll courier it anywhere in India",
    20,
    WHITE,
)
add_notes(
    s3,
    "BagDrop partners with local shops. Contactless OTP/NFC access. Unique delivery service - NO competitor offers this.",
)

# ============================================
# SLIDE 4: WHY GUJARAT
# ============================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, WHITE)
add_rect(s4, Inches(0), Inches(0), Inches(0.5), Inches(7.5), YELLOW)
add_heading(
    s4, Inches(1), Inches(0.5), Inches(11), Inches(1), "WHY GUJARAT?", 64, BLACK
)

advantages = [
    ("TOURIST HUB", "Rann Utsav: 5 Lakh+ visitors annually"),
    ("TRANSIT POINTS", "Bhuj Railway, Gandhidham Junction"),
    ("GROWTH", "Gujarat tourism growing at 12% CAGR"),
    ("STRATEGIC", "Gateway to Rann, near Mundra Port"),
    ("CULTURE", "Dholavira, handicrafts, year-round tourism"),
]

for i, (title, desc) in enumerate(advantages):
    y = Inches(1.8 + i * 0.9)
    add_rect(s4, Inches(1), y, Inches(0.15), Inches(0.5), BLACK)
    add_text(s4, Inches(1.3), y, Inches(3.5), Inches(0.4), title, 18, BLACK, bold=True)
    add_text(s4, Inches(5), y, Inches(7), Inches(0.4), desc, 16, DARK_GRAY)

add_rect(s4, Inches(8), Inches(5.5), Inches(4.8), Inches(1.5), BLACK)
add_text(
    s4,
    Inches(8.2),
    Inches(5.7),
    Inches(4.4),
    Inches(0.7),
    "5.5 CRORE",
    36,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_text(
    s4,
    Inches(8.2),
    Inches(6.4),
    Inches(4.4),
    Inches(0.4),
    "Annual Gujarat Tourists",
    14,
    WHITE,
    align=PP_ALIGN.CENTER,
)
add_notes(
    s4,
    "Kutch is perfect for launch. Rann Utsav brings 5L+ visitors. Major railway stations. 12% tourism growth.",
)

# ============================================
# SLIDE 5: TARGET MARKET - IMPROVED
# ============================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, WHITE)
add_rect(s5, Inches(0), Inches(0), Inches(0.5), Inches(7.5), YELLOW)
add_heading(
    s5, Inches(1), Inches(0.5), Inches(11), Inches(1), "TARGET MARKET", 56, BLACK
)

# Better structured segments
segments = [
    ("TOURISTS & BACKPACKERS", "Rann Utsav visitors, sightseers, 18-45 years"),
    ("TRAIN & BUS TRAVELERS", "Inter-city commuters, layover travelers"),
    ("STUDENTS", "Budget-conscious, exam/interview travel"),
    ("BUSINESS TRAVELERS", "Premium service, NFC preferred, 28-50 years"),
    ("FORGOTTEN LUGGAGE", "High-value delivery service, all ages"),
]

for i, (segment, desc) in enumerate(segments):
    y = Inches(1.6 + i * 1.1)
    color = YELLOW if i % 2 == 0 else BLACK
    txt_color = BLACK if i % 2 == 0 else WHITE

    add_rect(s5, Inches(1), y, Inches(5.5), Inches(0.95), color)
    add_text(
        s5,
        Inches(1.2),
        y + Inches(0.25),
        Inches(5.1),
        Inches(0.5),
        segment,
        18,
        txt_color,
        bold=True,
    )

    add_rect(s5, Inches(6.5), y, Inches(6.3), Inches(0.95), RGBColor(245, 245, 245))
    add_text(
        s5, Inches(6.7), y + Inches(0.25), Inches(5.9), Inches(0.5), desc, 16, DARK_GRAY
    )

add_notes(
    s5,
    "5 customer segments: Tourists, train/bus travelers, students, business travelers, and forgotten luggage users - our unique segment.",
)

# ============================================
# SLIDE 6: MARKET SIZE
# ============================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6, BLACK)
add_heading(
    s6,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.8),
    "MARKET OPPORTUNITY",
    56,
    YELLOW,
)

stats = [
    ("TAM", "₹2,000+ Cr", "Total Addressable"),
    ("SAM", "₹500 Cr", "Serviceable Market"),
    ("SOM", "₹36 Lakh", "Year 1 Target"),
]

for i, (label, value, desc) in enumerate(stats):
    x = Inches(0.5 + i * 4.2)
    color = YELLOW if i == 2 else RGBColor(40, 40, 40)
    txt_color = BLACK if i == 2 else YELLOW

    add_rect(s6, x, Inches(1.5), Inches(4), Inches(2), color)
    add_text(
        s6,
        x,
        Inches(1.7),
        Inches(4),
        Inches(0.8),
        value,
        36,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s6,
        x,
        Inches(2.5),
        Inches(4),
        Inches(0.4),
        label,
        20,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s6,
        x,
        Inches(3),
        Inches(4),
        Inches(0.4),
        desc,
        14,
        LIGHT_GRAY if i < 2 else DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )

add_text(
    s6,
    Inches(0.8),
    Inches(4.2),
    Inches(5.5),
    Inches(2),
    "GUJARAT:\n• 5.5 Cr domestic tourists\n• 5.2 Lakh foreign tourists\n• 12% CAGR growth",
    18,
    WHITE,
)
add_text(
    s6,
    Inches(7),
    Inches(4.2),
    Inches(5.5),
    Inches(2),
    "INDIA:\n• 2,400 Cr railway passengers\n• 15 Cr air passengers\n• ₹15 Lakh Cr tourism GDP",
    18,
    WHITE,
)
add_notes(
    s6,
    "TAM is ₹2,000+ Crore. Our Year 1 target: ₹36 Lakh, growing to ₹1.27 Crore by Year 2.",
)

# ============================================
# SLIDE 7: BUSINESS MODEL - IMPROVED TABLE FONTS
# ============================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7, WHITE)
add_heading(
    s7, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), "BUSINESS MODEL", 56, BLACK
)

# Bigger table with larger fonts
revenue = [
    ("STORAGE FEES", "30%", "₹50-200/booking"),
    ("DELIVERY FEES", "20%", "₹300-1,500"),
    ("NFC TOKENS", "100%", "₹199 one-time"),
    ("KIOSK ADS", "100%", "₹5,000/month"),
    ("PARTNER SUBS", "100%", "₹500/month"),
]

# Headers - BIGGER
headers = ["REVENUE STREAM", "MARGIN", "PRICING"]
for col, header in enumerate(headers):
    x = Inches(0.5 + col * 4.1)
    add_rect(s7, x, Inches(1.4), Inches(4.1), Inches(0.5), BLACK)
    add_text(
        s7,
        x,
        Inches(1.48),
        Inches(4.1),
        Inches(0.4),
        header,
        16,
        YELLOW,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

# Data - BIGGER FONTS
for row_idx, (stream, margin, price) in enumerate(revenue):
    y = Inches(1.9 + row_idx * 0.65)
    data = [stream, margin, price]
    for col, val in enumerate(data):
        x = Inches(0.5 + col * 4.1)
        bg_color = RGBColor(245, 245, 245) if row_idx % 2 == 0 else WHITE
        add_rect(s7, x, y, Inches(4.1), Inches(0.6), bg_color)
        txt_color = BLACK if col == 0 else (GREEN if col == 1 else DARK_GRAY)
        add_text(
            s7,
            x + Inches(0.15),
            y + Inches(0.15),
            Inches(3.8),
            Inches(0.35),
            val,
            15,
            txt_color,
            bold=(col == 0),
        )

# Pricing highlight
add_rect(s7, Inches(0.5), Inches(5.5), Inches(6), Inches(1.5), YELLOW)
add_text(
    s7,
    Inches(0.7),
    Inches(5.7),
    Inches(5.6),
    Inches(0.5),
    "STARTING AT ₹50",
    28,
    BLACK,
    bold=True,
)
add_text(
    s7,
    Inches(0.7),
    Inches(6.3),
    Inches(5.6),
    Inches(0.4),
    "For 6 hours - Cheaper than railway lockers",
    16,
    DARK_GRAY,
)
add_notes(
    s7,
    "5 revenue streams. Primary: 30% on storage. Starting price ₹50 for 6 hours beats railway lockers.",
)

# ============================================
# SLIDE 8: HOW IT WORKS
# ============================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8, WHITE)
add_heading(
    s8, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), "HOW IT WORKS", 56, BLACK
)
add_text(
    s8,
    Inches(0.8),
    Inches(1.1),
    Inches(11),
    Inches(0.5),
    "5 simple steps to freedom",
    22,
    LIGHT_GRAY,
)

steps = [
    ("01", "FIND", "Open app\nLocate point"),
    ("02", "BOOK", "Select slot\nPay via UPI"),
    ("03", "DROP", "Leave bags\nGet OTP"),
    ("04", "RETRIEVE", "Show token\nCollect bags"),
    ("05", "DELIVER", "Forgot bag?\nWe courier it!"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.3 + i * 2.6)
    color = YELLOW if i % 2 == 0 else BLACK
    txt_color = BLACK if i % 2 == 0 else YELLOW
    desc_color = DARK_GRAY if i % 2 == 0 else WHITE

    add_rect(s8, x, Inches(2), Inches(2.4), Inches(3.8), color)
    add_text(
        s8,
        x,
        Inches(2.2),
        Inches(2.4),
        Inches(0.6),
        num,
        36,
        txt_color,
        bold=True,
        font=HEADING_FONT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s8,
        x,
        Inches(2.9),
        Inches(2.4),
        Inches(0.5),
        title,
        22,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s8,
        x,
        Inches(3.5),
        Inches(2.4),
        Inches(1.8),
        desc,
        16,
        desc_color,
        align=PP_ALIGN.CENTER,
    )

add_rect(s8, Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.6), BLACK)
add_text(
    s8,
    Inches(0.8),
    Inches(6.3),
    Inches(11.8),
    Inches(0.5),
    "Average drop time: Under 2 minutes",
    20,
    YELLOW,
    bold=True,
)
add_notes(
    s8,
    "Simple 5-step process. Find, Book, Drop, Retrieve, Deliver. Under 2 minutes to drop bags.",
)

# ============================================
# SLIDE 9: TECHNOLOGY
# ============================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9, BLACK)
add_heading(
    s9, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "TECHNOLOGY", 64, YELLOW
)
add_text(
    s9,
    Inches(0.8),
    Inches(1.4),
    Inches(11),
    Inches(0.5),
    "Built for speed, security & scale",
    24,
    WHITE,
)

columns = [
    (
        "PAYMENTS",
        "• UPI (GPay, PhonePe)\n• Credit & Debit Cards\n• NFC Contactless\n• Wallet Integration",
    ),
    (
        "SECURITY",
        "• OTP Verification\n• NFC Token Access\n• End-to-End Encryption\n• Verified Partners",
    ),
    (
        "PLATFORM",
        "• iOS & Android App\n• Self-Service Kiosks\n• Partner Dashboard\n• Real-time Analytics",
    ),
]

for i, (title, content) in enumerate(columns):
    x = Inches(0.5 + i * 4.2)
    add_rect(s9, x, Inches(2.3), Inches(3.9), Inches(4), RGBColor(30, 30, 30))
    add_rect(s9, x, Inches(2.3), Inches(3.9), Inches(0.15), YELLOW)
    add_text(
        s9,
        x + Inches(0.2),
        Inches(2.6),
        Inches(3.5),
        Inches(0.5),
        title,
        22,
        YELLOW,
        bold=True,
    )
    add_text(
        s9, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(2.8), content, 16, WHITE
    )
add_notes(
    s9,
    "Tech stack: UPI/cards/NFC payments. OTP and NFC security. iOS/Android app, kiosks, partner dashboard.",
)

# ============================================
# SLIDE 10: COMPETITIVE ADVANTAGE - IMPROVED TABLE
# ============================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10, WHITE)
add_heading(
    s10,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.8),
    "WHY BAGDROP WINS",
    56,
    BLACK,
)

# Bigger table headers
add_rect(s10, Inches(0.5), Inches(1.4), Inches(4), Inches(0.5), BLACK)
add_text(
    s10,
    Inches(0.5),
    Inches(1.48),
    Inches(4),
    Inches(0.4),
    "FEATURE",
    16,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)

add_rect(s10, Inches(4.5), Inches(1.4), Inches(4), Inches(0.5), BLACK)
add_text(
    s10,
    Inches(4.5),
    Inches(1.48),
    Inches(4),
    Inches(0.4),
    "BAGDROP",
    16,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)

add_rect(s10, Inches(8.5), Inches(1.4), Inches(4.3), Inches(0.5), LIGHT_GRAY)
add_text(
    s10,
    Inches(8.5),
    Inches(1.48),
    Inches(4.3),
    Inches(0.4),
    "COMPETITORS",
    16,
    WHITE,
    bold=True,
    align=PP_ALIGN.CENTER,
)

comparisons = [
    ("Pricing", "₹50 for 6 hours", "₹100+/day"),
    ("Security", "OTP / NFC Tokens", "Physical Keys"),
    ("Payments", "UPI / Card / NFC", "Cash Only"),
    ("Delivery", "YES - Nationwide", "NO"),
    ("Booking", "App-based", "In-person"),
    ("Support", "24/7 Available", "Limited Hours"),
]

for row_idx, (feature, bagdrop, others) in enumerate(comparisons):
    y = Inches(1.9 + row_idx * 0.6)
    bg = RGBColor(245, 245, 245) if row_idx % 2 == 0 else WHITE

    add_rect(s10, Inches(0.5), y, Inches(4), Inches(0.55), bg)
    add_text(
        s10,
        Inches(0.7),
        y + Inches(0.12),
        Inches(3.6),
        Inches(0.35),
        feature,
        15,
        DARK_GRAY,
        bold=True,
    )

    add_rect(s10, Inches(4.5), y, Inches(4), Inches(0.55), bg)
    txt_color = GREEN if row_idx == 3 else BLACK
    add_text(
        s10,
        Inches(4.5),
        y + Inches(0.12),
        Inches(4),
        Inches(0.35),
        bagdrop,
        15,
        txt_color,
        align=PP_ALIGN.CENTER,
    )

    add_rect(s10, Inches(8.5), y, Inches(4.3), Inches(0.55), bg)
    add_text(
        s10,
        Inches(8.5),
        y + Inches(0.12),
        Inches(4.3),
        Inches(0.35),
        others,
        15,
        LIGHT_GRAY,
        align=PP_ALIGN.CENTER,
    )

# USP
add_rect(s10, Inches(0.5), Inches(5.7), Inches(12.333), Inches(0.8), BLACK)
add_text(
    s10,
    Inches(1),
    Inches(5.85),
    Inches(11.333),
    Inches(0.5),
    "UNIQUE USP: Forgotten luggage delivery - NO COMPETITOR offers this!",
    18,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_notes(
    s10,
    "We beat competitors on every metric. Our USP: forgotten luggage delivery service.",
)

# ============================================
# SLIDE 11: FINANCIAL PROJECTIONS - IMPROVED
# ============================================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s11, WHITE)
add_heading(
    s11,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.8),
    "FINANCIAL PROJECTIONS",
    48,
    BLACK,
)

# 3-year summary cards
years = [
    ("YEAR 1", "₹36.8L", "1.32L bookings"),
    ("YEAR 2", "₹1.27 Cr", "5L bookings"),
    ("YEAR 3", "₹3.82 Cr", "15L bookings"),
]

for i, (year, revenue, bookings) in enumerate(years):
    x = Inches(0.5 + i * 4.2)
    color = BLACK if i < 2 else YELLOW
    txt_color = YELLOW if i < 2 else BLACK

    add_rect(s11, x, Inches(1.3), Inches(4), Inches(1.6), color)
    add_text(
        s11,
        x,
        Inches(1.4),
        Inches(4),
        Inches(0.5),
        year,
        20,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s11,
        x,
        Inches(1.9),
        Inches(4),
        Inches(0.6),
        revenue,
        32,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s11,
        x,
        Inches(2.5),
        Inches(4),
        Inches(0.35),
        bookings,
        14,
        LIGHT_GRAY if i < 2 else DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )

# Monthly growth - simplified
add_text(
    s11,
    Inches(0.8),
    Inches(3.2),
    Inches(11),
    Inches(0.4),
    "YEAR 1 MONTHLY GROWTH",
    18,
    BLACK,
    bold=True,
)

months = [
    ("M1", "500", "₹9K"),
    ("M3", "2K", "₹39K"),
    ("M6", "7.5K", "₹1.57L"),
    ("M9", "15K", "₹3.6L"),
    ("M12", "30K", "₹7.65L"),
]

for i, (m, b, r) in enumerate(months):
    x = Inches(0.5 + i * 2.5)
    color = YELLOW if i < 4 else BLACK
    txt_color = BLACK if i < 4 else YELLOW

    add_rect(s11, x, Inches(3.7), Inches(2.3), Inches(1.8), color)
    add_text(
        s11,
        x,
        Inches(3.85),
        Inches(2.3),
        Inches(0.4),
        m,
        22,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s11,
        x,
        Inches(4.3),
        Inches(2.3),
        Inches(0.35),
        b + " bookings",
        12,
        txt_color,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s11,
        x,
        Inches(4.7),
        Inches(2.3),
        Inches(0.5),
        r,
        18,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

# Break-even
add_rect(s11, Inches(0.5), Inches(5.8), Inches(6), Inches(0.8), GREEN)
add_text(
    s11,
    Inches(0.7),
    Inches(5.95),
    Inches(5.6),
    Inches(0.5),
    "BREAK-EVEN: Month 10-11 | Profitable from Q4",
    16,
    WHITE,
    bold=True,
)

# Team growth
add_rect(s11, Inches(7), Inches(5.8), Inches(5.8), Inches(0.8), BLACK)
add_text(
    s11,
    Inches(7.2),
    Inches(5.95),
    Inches(5.4),
    Inches(0.5),
    "TEAM: 8 → 25 → 60 people (Y1 → Y3)",
    16,
    YELLOW,
)
add_notes(
    s11,
    "Year 1: ₹36.8L. Year 3: ₹3.82Cr. Break-even at Month 10-11. Profitable from Q4.",
)

# ============================================
# SLIDE 12: THE ASK
# ============================================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s12, BLACK)
add_heading(
    s12, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "THE ASK", 64, YELLOW
)

add_heading(
    s12, Inches(0.8), Inches(1.5), Inches(6), Inches(1.5), "₹36 LAKH", 72, YELLOW
)
add_text(
    s12,
    Inches(0.8),
    Inches(3),
    Inches(6),
    Inches(0.5),
    "SEED FUNDING (~$43K USD)",
    24,
    WHITE,
)

breakdown = [
    ("Tech Development", "₹5 Lakh"),
    ("6-Month Runway", "₹18 Lakh"),
    ("Marketing", "₹6.5 Lakh"),
    ("Equipment", "₹3 Lakh"),
    ("Contingency", "₹3.25 Lakh"),
]

for i, (item, amt) in enumerate(breakdown):
    y = Inches(1.5 + i * 0.7)
    add_text(s12, Inches(7), y, Inches(3.5), Inches(0.4), item, 16, WHITE)
    add_text(s12, Inches(10.8), y, Inches(2), Inches(0.4), amt, 16, YELLOW, bold=True)

add_rect(s12, Inches(0.5), Inches(4.5), Inches(5.5), Inches(1.5), YELLOW)
add_text(
    s12,
    Inches(0.7),
    Inches(4.7),
    Inches(5.1),
    Inches(0.5),
    "EXPECTED RETURNS",
    18,
    BLACK,
    bold=True,
)
add_text(
    s12,
    Inches(0.7),
    Inches(5.3),
    Inches(5.1),
    Inches(0.5),
    "5-10x for seed investors",
    24,
    BLACK,
    bold=True,
)

add_rect(s12, Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.6), YELLOW)
add_text(
    s12,
    Inches(1),
    Inches(6.35),
    Inches(11.333),
    Inches(0.4),
    "Target Valuation: ₹50+ Crore by Year 5",
    18,
    BLACK,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_notes(
    s12,
    "₹36 Lakh seed funding. Expected 5-10x returns. Target: ₹50+ Cr valuation by Year 5.",
)

# ============================================
# SLIDE 13: TEAM - IMPROVED LAYOUT
# ============================================
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s13, WHITE)
add_rect(s13, Inches(0), Inches(0), Inches(0.5), Inches(7.5), YELLOW)
add_heading(
    s13, Inches(1), Inches(0.5), Inches(11), Inches(0.8), "TEAM STRUCTURE", 56, BLACK
)

team = [
    ("CEO / FOUNDER", "Strategy, vision, investor relations"),
    ("CTO", "App development, tech infrastructure"),
    ("OPERATIONS MGR", "Partner onboarding, quality control"),
    ("LOGISTICS", "Delivery partners, courier management"),
    ("SUPPORT LEAD", "24/7 phone & app support"),
]

for i, (role, desc) in enumerate(team):
    y = Inches(1.6 + i * 1.0)

    add_rect(s13, Inches(1), y, Inches(4), Inches(0.85), BLACK)
    add_text(
        s13,
        Inches(1.2),
        y + Inches(0.25),
        Inches(3.6),
        Inches(0.4),
        role,
        16,
        YELLOW,
        bold=True,
    )

    add_rect(s13, Inches(5), y, Inches(7.5), Inches(0.85), RGBColor(245, 245, 245))
    add_text(
        s13,
        Inches(5.2),
        y + Inches(0.25),
        Inches(7.1),
        Inches(0.4),
        desc,
        15,
        DARK_GRAY,
    )

# Stats at bottom
add_rect(s13, Inches(1), Inches(6.8), Inches(3.5), Inches(0.5), YELLOW)
add_text(
    s13,
    Inches(1),
    Inches(6.85),
    Inches(3.5),
    Inches(0.4),
    "8 Team Members by Year 1",
    14,
    BLACK,
    bold=True,
    align=PP_ALIGN.CENTER,
)

add_rect(s13, Inches(5), Inches(6.8), Inches(3.5), Inches(0.5), BLACK)
add_text(
    s13,
    Inches(5),
    Inches(6.85),
    Inches(3.5),
    Inches(0.4),
    "HQ: Kutch, Gujarat",
    14,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_notes(s13, "5 core roles. Growing to 8 members by Year 1. Based in Kutch, Gujarat.")

# ============================================
# SLIDE 14: MARKETING - RESTRUCTURED
# ============================================
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s14, WHITE)
add_heading(
    s14,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.8),
    "MARKETING STRATEGY",
    48,
    BLACK,
)

# Top row - 3 cards
channels_top = [
    ("DIGITAL ADS", "₹55K/mo", "Instagram, Facebook, Google"),
    ("INFLUENCERS", "₹30K/mo", "Travel vloggers, creators"),
    ("TRANSIT HUBS", "₹25K/mo", "Stations, airports, flyers"),
]

for i, (title, budget, desc) in enumerate(channels_top):
    x = Inches(0.5 + i * 4.2)
    color = YELLOW if i != 1 else BLACK
    txt_color = BLACK if i != 1 else YELLOW

    add_rect(s14, x, Inches(1.4), Inches(4), Inches(2), color)
    add_text(
        s14,
        x + Inches(0.2),
        Inches(1.6),
        Inches(3.6),
        Inches(0.4),
        title,
        18,
        txt_color,
        bold=True,
    )
    add_text(
        s14,
        x + Inches(0.2),
        Inches(2.1),
        Inches(3.6),
        Inches(0.5),
        desc,
        14,
        DARK_GRAY if i != 1 else WHITE,
    )
    add_text(
        s14,
        x + Inches(0.2),
        Inches(2.8),
        Inches(3.6),
        Inches(0.4),
        budget,
        16,
        GREEN if i != 1 else YELLOW,
        bold=True,
    )

# Bottom row - 2 cards
add_rect(s14, Inches(0.5), Inches(3.7), Inches(6), Inches(1.6), BLACK)
add_text(
    s14,
    Inches(0.7),
    Inches(3.9),
    Inches(5.6),
    Inches(0.4),
    "REFERRAL PROGRAM",
    18,
    YELLOW,
    bold=True,
)
add_text(
    s14,
    Inches(0.7),
    Inches(4.4),
    Inches(5.6),
    Inches(0.8),
    "₹50 credit for every friend who signs up\nStackable rewards, no expiry for 6 months",
    14,
    WHITE,
)

add_rect(s14, Inches(6.8), Inches(3.7), Inches(6), Inches(1.6), YELLOW)
add_text(
    s14,
    Inches(7),
    Inches(3.9),
    Inches(5.6),
    Inches(0.4),
    "PARTNER ONBOARDING",
    18,
    BLACK,
    bold=True,
)
add_text(
    s14,
    Inches(7),
    Inches(4.4),
    Inches(5.6),
    Inches(0.8),
    "70% revenue share to shop owners\nNo lock-in, weekly payouts",
    14,
    DARK_GRAY,
)

# Budget summary
add_rect(s14, Inches(0.5), Inches(5.6), Inches(12.333), Inches(0.7), BLACK)
add_text(
    s14,
    Inches(1),
    Inches(5.75),
    Inches(11.333),
    Inches(0.4),
    "ANNUAL MARKETING BUDGET: ₹6.45 LAKH",
    18,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)

# KPIs
add_text(
    s14,
    Inches(0.8),
    Inches(6.5),
    Inches(11),
    Inches(0.5),
    "TARGETS: 50K Downloads | 100 Partners | 30K Monthly Bookings (Year 1)",
    14,
    DARK_GRAY,
)
add_notes(
    s14,
    "Multi-channel: Digital (₹55K), influencers (₹30K), transit hubs (₹25K). Referrals + partner onboarding. Total: ₹6.45L annually.",
)

print("Slides 1-14 complete")
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V5.pptx"
)
print("Saved slides 1-14 to V5")
