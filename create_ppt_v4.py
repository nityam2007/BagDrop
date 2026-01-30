#!/usr/bin/env python3
"""
BagDrop Presentation V4 - Professional 20-Slide Deck
With Speaker Notes, Tables, and Graphics
Team: TeamBack | Kutch, Gujarat
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand Colors
YELLOW = RGBColor(255, 215, 0)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(245, 245, 245)


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK
    if subtitle:
        tf.add_paragraph().text = subtitle
        tf.paragraphs[1].font.size = Pt(16)
        tf.paragraphs[1].font.color.rgb = GRAY


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_bullet_list(slide, items, left, top, width, font_size=14):
    box = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.4 + 0.3))
    tf = box.text_frame
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.level = 0


# ===================== SLIDE 1: TITLE =====================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
box = s1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
tf = box.text_frame
tf.text = "BagDrop"
p = tf.paragraphs[0]
p.font.size = Pt(80)
p.font.bold = True
p.font.color.rgb = BLACK
p.alignment = PP_ALIGN.CENTER

tf.add_paragraph().text = "India's Smart Luggage Storage Network"
tf.paragraphs[1].font.size = Pt(28)
tf.paragraphs[1].font.color.rgb = GRAY
tf.paragraphs[1].alignment = PP_ALIGN.CENTER

tf.add_paragraph().text = ""
tf.add_paragraph().text = "TeamBack | Kutch, Gujarat | January 2026"
tf.paragraphs[3].font.size = Pt(16)
tf.paragraphs[3].font.color.rgb = DARK_GRAY
tf.paragraphs[3].alignment = PP_ALIGN.CENTER

accent = s1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.5), Inches(5.333), Inches(0.15)
)
accent.fill.solid()
accent.fill.fore_color.rgb = YELLOW
accent.line.fill.background()

add_notes(
    s1,
    "Welcome judges! We are TeamBack from Kutch, Gujarat. Today we present BagDrop - India's smart luggage storage network. 70% of travelers face luggage problems. We're solving this with technology. This is a $15 Lakh Crore market opportunity. Let's dive in.",
)

# ===================== SLIDE 2: PROBLEM =====================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)
add_title(s2, "The Problem", "70% of travelers face luggage issues")

# Quote box
quote_box = s2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(6), Inches(2.2)
)
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
quote_box.line.color.rgb = RGBColor(200, 50, 50)

txt = s2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.4), Inches(1.8))
tf = txt.text_frame
tf.text = '"Checked out at 10 AM, train at 7 PM..."'
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = RGBColor(150, 0, 0)

problems = [
    "• Heavy bags ruin sightseeing plans",
    "• Railway lockers cost ₹100+/day",
    "• Often full or broken",
    "• No solution for forgotten luggage",
]
add_bullet_list(s2, problems, Inches(0.8), Inches(2.8), Inches(5.4))

# Stats
stats = [
    ("5.5 Crore", "Annual tourists in Gujarat"),
    ("2,400 Crore", "Railway passengers India"),
    ("70%", "Face luggage issues"),
]
for i, (stat, label) in enumerate(stats):
    y = Inches(1.5 + i * 1.0)
    box = s2.shapes.add_textbox(Inches(7.5), y, Inches(2), Inches(0.5))
    tf = box.text_frame
    tf.text = stat
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = YELLOW

    box2 = s2.shapes.add_textbox(Inches(9.7), y + Inches(0.1), Inches(3), Inches(0.4))
    tf2 = box2.text_frame
    tf2.text = label
    tf2.paragraphs[0].font.size = Pt(13)
    tf2.paragraphs[0].font.color.rgb = DARK_GRAY

add_notes(
    s2,
    "Paint the picture: Early checkout, late departure. You're dragging 20kg bags through temples, cafes, crowded markets. Railway lockers are expensive and unreliable. 5.5 Crore tourists visit Gujarat annually - imagine their frustration. This is a universal pain point with no good solution.",
)

# ===================== SLIDE 3: SOLUTION =====================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3)
add_title(s3, "Our Solution: BagDrop", "Store. Explore. Relax.")

steps = [
    ("1", "FIND", "Locate nearby BagDrop"),
    ("2", "BOOK", "Pay via UPI/Card"),
    ("3", "DROP", "Get OTP/NFC token"),
    ("4", "RETRIEVE", "Show token, collect"),
    ("5", "DELIVER", "We courier forgotten bags"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(1.4 + (i // 3) * 2.2)

    circle = s3.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = YELLOW
    circle.line.fill.background()

    num_box = s3.shapes.add_textbox(x, y + Inches(0.05), Inches(0.5), Inches(0.4))
    num_box.text_frame.text = num
    num_box.text_frame.paragraphs[0].font.size = Pt(16)
    num_box.text_frame.paragraphs[0].font.bold = True
    num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    title_box = s3.shapes.add_textbox(x + Inches(0.7), y, Inches(3.3), Inches(0.4))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(16)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    desc_box = s3.shapes.add_textbox(
        x + Inches(0.7), y + Inches(0.4), Inches(3.3), Inches(0.8)
    )
    desc_box.text_frame.text = desc
    desc_box.text_frame.paragraphs[0].font.size = Pt(12)
    desc_box.text_frame.paragraphs[0].font.color.rgb = GRAY

# USP banner
banner = s3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2)
)
banner.fill.solid()
banner.fill.fore_color.rgb = BLACK
banner.line.fill.background()

banner_txt = s3.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(0.8))
banner_tf = banner_txt.text_frame
banner_tf.text = "Under 2 min drop time | Starting ₹50/6hrs | 24/7 Support | Forgotten luggage delivery"
banner_tf.paragraphs[0].font.size = Pt(20)
banner_tf.paragraphs[0].font.color.rgb = YELLOW
banner_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_notes(
    s3,
    "Here's how BagDrop works. The key differentiator: we deliver forgotten luggage anywhere in India. No competitor does this. Under 2 minutes to drop bags. Starting at just ₹50 for 6 hours - cheaper than railway lockers. 24/7 support. Emphasize speed, affordability, and the unique delivery feature.",
)

# ===================== SLIDE 4: WHY GUJARAT =====================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4)
add_title(s4, "Why Kutch, Gujarat?", "Strategic beachhead market")

# Create table manually
table_data = [
    ["Factor", "Advantage"],
    ["Tourist Hub", "Rann Utsav: 5L+ visitors annually"],
    ["Transit Points", "Bhuj Railway, Gandhidham Junction"],
    ["Growth Rate", "Gujarat tourism: 12% CAGR"],
    ["Strategic", "Gateway to Rann, near Mundra Port"],
    ["Cultural", "Dholavira, handicrafts, year-round"],
]

for row_idx, row in enumerate(table_data):
    for col_idx, cell_text in enumerate(row):
        x = Inches(0.5 + col_idx * 6)
        y = Inches(1.4 + row_idx * 0.7)
        w = Inches(6)
        h = Inches(0.7)

        cell = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = YELLOW
            cell.line.color.rgb = YELLOW
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 1 else LIGHT_GRAY
            cell.line.color.rgb = GRAY

        txt = s4.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), h - Inches(0.3)
        )
        txt.text_frame.text = cell_text
        txt.text_frame.paragraphs[0].font.size = Pt(13 if row_idx > 0 else 14)
        txt.text_frame.paragraphs[0].font.bold = row_idx == 0
        txt.text_frame.paragraphs[0].font.color.rgb = BLACK

highlight = s4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(5), Inches(1.5)
)
highlight.fill.solid()
highlight.fill.fore_color.rgb = RGBColor(255, 250, 220)
highlight.line.color.rgb = YELLOW

hl_txt = s4.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(4.6), Inches(1.1))
hl_tf = hl_txt.text_frame
hl_tf.text = "Rann Utsav Season (Nov-Feb)"
hl_tf.add_paragraph().text = "Peak: 5-7 Lakh visitors"
hl_tf.add_paragraph().text = "Perfect launch timing!"
hl_tf.paragraphs[0].font.size = Pt(16)
hl_tf.paragraphs[0].font.bold = True
hl_tf.paragraphs[0].font.color.rgb = BLACK
for p in hl_tf.paragraphs[1:]:
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

add_notes(
    s4,
    "Why Kutch? Rann Utsav alone brings 5-7 Lakh visitors in 4 months. Major railway stations at Bhuj and Gandhidham. Gujarat tourism growing at 12% CAGR. We have the perfect test market with high tourist traffic and limited competition. This is our beachhead before pan-India expansion.",
)

print("Slides 1-4 complete")

# ===================== SLIDE 5: TARGET MARKET =====================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5)
add_title(s5, "Target Market", "5 primary customer segments")

segments = [
    ("Tourists & Backpackers", "18-45 years", "Rann Utsav, sightseeing"),
    ("Train & Bus Travelers", "20-55 years", "Inter-city commuters"),
    ("Students", "17-25 years", "Exams, interviews, budget"),
    ("Business Travelers", "28-50 years", "Meetings, premium service"),
    ("Forgotten Luggage", "All ages", "High-value delivery need"),
]

for i, (segment, demo, use) in enumerate(segments):
    y = Inches(1.4 + i * 1.0)

    # Segment name
    box = s5.shapes.add_textbox(Inches(0.5), y, Inches(4), Inches(0.5))
    box.text_frame.text = segment
    box.text_frame.paragraphs[0].font.size = Pt(16)
    box.text_frame.paragraphs[0].font.bold = True
    box.text_frame.paragraphs[0].font.color.rgb = BLACK

    # Demographics
    box2 = s5.shapes.add_textbox(Inches(5), y, Inches(3), Inches(0.4))
    box2.text_frame.text = demo
    box2.text_frame.paragraphs[0].font.size = Pt(12)
    box2.text_frame.paragraphs[0].font.color.rgb = GRAY

    # Use case
    box3 = s5.shapes.add_textbox(Inches(8), y, Inches(4.5), Inches(0.4))
    box3.text_frame.text = use
    box3.text_frame.paragraphs[0].font.size = Pt(12)
    box3.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

add_notes(
    s5,
    "We serve 5 key segments. Tourists visiting Rann Utsav. Train travelers with long layovers. Students moving between cities. Business travelers needing quick storage. And our unique segment - people who forgot their luggage and need it delivered. Each has different needs but shares the same pain point.",
)

# ===================== SLIDE 6: MARKET SIZE =====================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6)
add_title(s6, "Market Opportunity", "Massive addressable market")

# Gujarat stats
box1 = s6.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(6), Inches(2.2)
)
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(240, 248, 255)
box1.line.color.rgb = RGBColor(100, 150, 200)

t1 = s6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.4), Inches(1.8))
tf1 = t1.text_frame
tf1.text = "Gujarat Market"
tf1.paragraphs[0].font.size = Pt(20)
tf1.paragraphs[0].font.bold = True
tf1.paragraphs[0].font.color.rgb = RGBColor(50, 100, 150)
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "5.5 Crore domestic tourists annually"
tf1.add_paragraph().text = "5.2 Lakh foreign tourists"
tf1.add_paragraph().text = "12% CAGR tourism growth"
for p in tf1.paragraphs[1:]:
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_GRAY

# India stats
box2 = s6.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.4), Inches(5.8), Inches(2.2)
)
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(255, 250, 240)
box2.line.color.rgb = YELLOW

t2 = s6.shapes.add_textbox(Inches(7.3), Inches(1.6), Inches(5.2), Inches(1.8))
tf2 = t2.text_frame
tf2.text = "India Market"
tf2.paragraphs[0].font.size = Pt(20)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = BLACK
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "2,400 Crore railway passengers"
tf2.add_paragraph().text = "15 Crore air passengers"
tf2.add_paragraph().text = "₹15 Lakh Crore tourism GDP"
for p in tf2.paragraphs[1:]:
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_GRAY

# TAM SAM SOM
metrics = [
    ("TAM", "Total Addressable", "₹2,000+ Crore", "All travelers needing storage"),
    ("SAM", "Serviceable", "₹500 Crore", "Gujarat + 10 major cities"),
    ("SOM", "Obtainable", "₹36 Lakh Y1", "Our Year 1 target"),
]

for i, (label, desc, value, note) in enumerate(metrics):
    y = Inches(4 + i * 0.9)

    l_box = s6.shapes.add_textbox(Inches(0.5), y, Inches(2.5), Inches(0.5))
    l_box.text_frame.text = label
    l_box.text_frame.paragraphs[0].font.size = Pt(18)
    l_box.text_frame.paragraphs[0].font.bold = True
    l_box.text_frame.paragraphs[0].font.color.rgb = YELLOW if i == 2 else BLACK

    d_box = s6.shapes.add_textbox(Inches(3), y, Inches(3), Inches(0.4))
    d_box.text_frame.text = desc
    d_box.text_frame.paragraphs[0].font.size = Pt(12)
    d_box.text_frame.paragraphs[0].font.color.rgb = GRAY

    v_box = s6.shapes.add_textbox(Inches(6.5), y, Inches(2.5), Inches(0.5))
    v_box.text_frame.text = value
    v_box.text_frame.paragraphs[0].font.size = Pt(16)
    v_box.text_frame.paragraphs[0].font.bold = True
    v_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    n_box = s6.shapes.add_textbox(Inches(9.5), y, Inches(3.3), Inches(0.4))
    n_box.text_frame.text = note
    n_box.text_frame.paragraphs[0].font.size = Pt(11)
    n_box.text_frame.paragraphs[0].font.color.rgb = GRAY

add_notes(
    s6,
    "The market is massive. 5.5 Crore tourists in Gujarat alone. 2,400 Crore railway passengers across India. TAM is ₹2,000+ Crore. Even capturing 1% is huge. We're targeting SAM of ₹500 Crore in Gujarat + major cities first. Our Year 1 goal is ₹36 Lakh revenue, growing to ₹1.27 Crore by Year 2.",
)

print("Slides 5-6 complete")

# Save progress
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4.pptx"
)
print("Saved slides 1-6")
