#!/usr/bin/env python3
"""
BagDrop Presentation V4 - REDESIGNED with Modern Corporate Theme
Design: Blue/White/Gray color scheme per prompt.md
Professional, clean layout with icons and visual placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# MODERN CORPORATE COLOR SCHEME (per prompt.md)
PRIMARY_BLUE = RGBColor(41, 98, 255)  # Modern blue
SECONDARY_BLUE = RGBColor(66, 133, 244)  # Light blue
DARK_BLUE = RGBColor(25, 55, 109)  # Dark blue for headers
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(248, 249, 250)  # Background gray
MEDIUM_GRAY = RGBColor(128, 128, 128)  # Text gray
DARK_GRAY = RGBColor(60, 64, 67)  # Dark text
ACCENT_GRAY = RGBColor(232, 234, 237)  # Border gray
SUCCESS_GREEN = RGBColor(52, 168, 83)  # For positive metrics


def add_background(slide, color=WHITE):
    """Add solid background"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_header_bar(slide):
    """Add blue header bar at top"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_BLUE
    bar.line.fill.background()


def add_footer_bar(slide):
    """Add gray footer bar at bottom"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GRAY
    bar.line.fill.background()

    # Add page number placeholder
    footer = slide.shapes.add_textbox(
        Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.2)
    )
    footer.text_frame.text = "BagDrop"
    footer.text_frame.paragraphs[0].font.size = Pt(10)
    footer.text_frame.paragraphs[0].font.color.rgb = MEDIUM_GRAY


def add_title_style(slide, title, subtitle="", has_icon=False):
    """Add styled title with blue accent"""
    # Blue accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.6)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY_BLUE
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.45), Inches(11.5), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Calibri"

    if subtitle:
        tf.add_paragraph().text = subtitle
        tf.paragraphs[1].font.size = Pt(14)
        tf.paragraphs[1].font.color.rgb = MEDIUM_GRAY
        tf.paragraphs[1].font.name = "Calibri"


def add_content_box(
    slide, left, top, width, height, title="", content=None, is_highlight=False
):
    """Add styled content box with optional title"""
    # Box background
    if is_highlight:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(232, 240, 254)  # Light blue highlight
        box.line.color.rgb = PRIMARY_BLUE
    else:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = ACCENT_GRAY

    current_top = top + Inches(0.15)

    # Title if provided
    if title:
        title_box = slide.shapes.add_textbox(
            left + Inches(0.2), current_top, width - Inches(0.4), Inches(0.35)
        )
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(14)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        current_top += Inches(0.4)

    # Content
    if content:
        content_box = slide.shapes.add_textbox(
            left + Inches(0.2), current_top, width - Inches(0.4), height - Inches(0.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        for item in content:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(6)


def add_icon_placeholder(slide, left, top, icon_type="image"):
    """Add icon placeholder"""
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.8), Inches(0.8)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = ACCENT_GRAY
    placeholder.line.color.rgb = MEDIUM_GRAY

    icon_text = slide.shapes.add_textbox(
        left, top + Inches(0.25), Inches(0.8), Inches(0.3)
    )
    icon_text.text_frame.text = "[Icon]"
    icon_text.text_frame.paragraphs[0].font.size = Pt(9)
    icon_text.text_frame.paragraphs[0].font.color.rgb = MEDIUM_GRAY
    icon_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_image_placeholder(slide, left, top, width, height, label="Image"):
    """Add image placeholder with label"""
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = ACCENT_GRAY
    placeholder.line.color.rgb = MEDIUM_GRAY
    placeholder.line.width = Pt(1)

    label_box = slide.shapes.add_textbox(
        left, top + height / 2 - Inches(0.15), width, Inches(0.3)
    )
    label_box.text_frame.text = f"[{label}]"
    label_box.text_frame.paragraphs[0].font.size = Pt(10)
    label_box.text_frame.paragraphs[0].font.color.rgb = MEDIUM_GRAY
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_stat_card(slide, left, top, stat, label, is_large=False):
    """Add statistic card"""
    width = Inches(3) if is_large else Inches(2.5)
    height = Inches(1.2) if is_large else Inches(1)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = PRIMARY_BLUE
    card.line.width = Pt(2)

    stat_box = slide.shapes.add_textbox(left, top + Inches(0.15), width, Inches(0.5))
    stat_box.text_frame.text = stat
    stat_box.text_frame.paragraphs[0].font.size = Pt(28) if is_large else Pt(24)
    stat_box.text_frame.paragraphs[0].font.bold = True
    stat_box.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    stat_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    label_box = slide.shapes.add_textbox(left, top + Inches(0.6), width, Inches(0.4))
    label_box.text_frame.text = label
    label_box.text_frame.paragraphs[0].font.size = Pt(10)
    label_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_process_step(slide, x, y, number, title, description):
    """Add numbered process step with circle"""
    # Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_BLUE
    circle.line.fill.background()

    # Number
    num_box = slide.shapes.add_textbox(x, y + Inches(0.08), Inches(0.5), Inches(0.35))
    num_box.text_frame.text = str(number)
    num_box.text_frame.paragraphs[0].font.size = Pt(18)
    num_box.text_frame.paragraphs[0].font.bold = True
    num_box.text_frame.paragraphs[0].font.color.rgb = WHITE
    num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(x + Inches(0.6), y, Inches(3), Inches(0.3))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(13)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

    # Description
    desc_box = slide.shapes.add_textbox(
        x + Inches(0.6), y + Inches(0.3), Inches(3), Inches(0.5)
    )
    desc_box.text_frame.text = description
    desc_box.text_frame.paragraphs[0].font.size = Pt(10)
    desc_box.text_frame.paragraphs[0].font.color.rgb = MEDIUM_GRAY


def add_notes(slide, text):
    """Add speaker notes"""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ==========================================
# SLIDE 1: TITLE / COVER
# ==========================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s1, WHITE)
add_header_bar(s1)
add_footer_bar(s1)

# Main title
main_box = s1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(1.2))
main_tf = main_box.text_frame
main_tf.text = "BagDrop"
main_tf.paragraphs[0].font.size = Pt(60)
main_tf.paragraphs[0].font.bold = True
main_tf.paragraphs[0].font.color.rgb = DARK_BLUE

main_tf.add_paragraph().text = "Smart Luggage Storage & Delivery"
main_tf.paragraphs[1].font.size = Pt(28)
main_tf.paragraphs[1].font.color.rgb = PRIMARY_BLUE

# Tagline
main_tf.add_paragraph().text = ""
main_tf.add_paragraph().text = "Secure • Contactless • Affordable"
main_tf.paragraphs[3].font.size = Pt(16)
main_tf.paragraphs[3].font.color.rgb = MEDIUM_GRAY

# Right side - image placeholders
add_image_placeholder(s1, Inches(9), Inches(1.5), Inches(3.5), Inches(2), "Logo")
add_image_placeholder(
    s1, Inches(9), Inches(4), Inches(3.5), Inches(2), "App Screenshot"
)

# Team info at bottom
team_box = s1.shapes.add_textbox(Inches(0.5), Inches(6), Inches(8), Inches(0.8))
team_tf = team_box.text_frame
team_tf.text = "TeamBack | Kutch, Gujarat"
team_tf.add_paragraph().text = "January 2026"
for p in team_tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = MEDIUM_GRAY

add_notes(
    s1,
    "Welcome! We're TeamBack from Kutch, Gujarat. BagDrop is India's smart luggage storage network solving a problem faced by 70% of travelers. Modern, secure, and affordable.",
)

# ==========================================
# SLIDE 2: PROBLEM STATEMENT
# ==========================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s2, WHITE)
add_header_bar(s2)
add_footer_bar(s2)
add_title_style(s2, "The Problem", "Travel storage challenges")

# Left side - problem points
add_content_box(
    s2,
    Inches(0.5),
    Inches(1.5),
    Inches(6),
    Inches(3.5),
    "Key Pain Points",
    [
        "✓ Early checkout, late departure (hours of waiting)",
        "✓ Carrying heavy bags ruins travel experience",
        "✓ Railway cloakrooms crowded, expensive (₹100+/day)",
        "✓ Forgotten luggage creates stress and loss",
    ],
)

# Right side - stats
add_stat_card(s2, Inches(7), Inches(1.5), "70%", "Travelers face luggage issues")
add_stat_card(s2, Inches(7), Inches(2.8), "5.5 Cr", "Annual Gujarat tourists")
add_stat_card(s2, Inches(7), Inches(4.1), "2,400 Cr", "India railway passengers")

# Image placeholder
add_image_placeholder(
    s2, Inches(10.5), Inches(1.5), Inches(2.3), Inches(2), "Travelers with Luggage"
)

add_notes(
    s2,
    "70% of travelers face luggage problems. Early checkout with late departure means hours dragging heavy bags. Railway lockers are expensive and unreliable. This is a universal pain point with no good solution.",
)

# ==========================================
# SLIDE 3: OUR SOLUTION
# ==========================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s3, WHITE)
add_header_bar(s3)
add_footer_bar(s3)
add_title_style(s3, "Our Solution", "Store bags anywhere, anytime")

# Solution features
features = [
    ("Partner Network", "Verified local shops as storage points"),
    ("Contactless Access", "OTP & NFC token-based retrieval"),
    ("Smart App", "Book, pay, and track in real-time"),
    ("Delivery Service", "Forgotten luggage courier anywhere"),
]

for i, (title, desc) in enumerate(features):
    y = Inches(1.5 + i * 1.1)
    add_icon_placeholder(s3, Inches(0.5), y)

    title_box = s3.shapes.add_textbox(Inches(1.5), y, Inches(4.5), Inches(0.3))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(14)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

    desc_box = s3.shapes.add_textbox(
        Inches(1.5), y + Inches(0.3), Inches(4.5), Inches(0.5)
    )
    desc_box.text_frame.text = desc
    desc_box.text_frame.paragraphs[0].font.size = Pt(11)
    desc_box.text_frame.paragraphs[0].font.color.rgb = MEDIUM_GRAY

# Right side - images
add_image_placeholder(
    s3, Inches(7), Inches(1.5), Inches(3), Inches(1.5), "App Interface"
)
add_image_placeholder(
    s3, Inches(10.5), Inches(1.5), Inches(2.3), Inches(1.5), "Storage Shop"
)
add_image_placeholder(s3, Inches(7), Inches(3.2), Inches(3), Inches(1.5), "Kiosk")
add_image_placeholder(
    s3, Inches(10.5), Inches(3.2), Inches(2.3), Inches(1.5), "NFC Token"
)

add_notes(
    s3,
    "BagDrop partners with local shops for secure storage. Contactless OTP/NFC access. Smart mobile app for booking. Unique delivery service for forgotten luggage - no competitor offers this.",
)

# ==========================================
# SLIDE 4: VALUE PROPOSITION
# ==========================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(s4, WHITE)
add_header_bar(s4)
add_footer_bar(s4)
add_title_style(s4, "Why BagDrop is Better", "Value proposition comparison")

# Comparison table header
headers = ["Feature", "BagDrop", "Traditional"]
for col, header in enumerate(headers):
    x = Inches(0.5 + col * 4)
    cell = s4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(4), Inches(0.4)
    )
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY_BLUE if col > 0 else DARK_BLUE
    cell.line.fill.background()

    txt = s4.shapes.add_textbox(x + Inches(0.2), Inches(1.55), Inches(3.6), Inches(0.3))
    txt.text_frame.text = header
    txt.text_frame.paragraphs[0].font.size = Pt(12)
    txt.text_frame.paragraphs[0].font.bold = True
    txt.text_frame.paragraphs[0].font.color.rgb = WHITE
    txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Comparison data
comparisons = [
    ("Pricing", "₹50 for 6 hrs", "₹100+/day"),
    ("Security", "OTP/NFC tokens", "Keys/Tags"),
    ("Payments", "UPI/Card/NFC", "Cash only"),
    ("Delivery", "✓ Available", "✗ Not available"),
    ("Booking", "App-based", "In-person only"),
    ("Support", "24/7 Phone/App", "Limited hours"),
]

for row_idx, (feature, bagdrop, trad) in enumerate(comparisons):
    y = Inches(1.9 + row_idx * 0.5)
    data = [feature, bagdrop, trad]
    for col, val in enumerate(data):
        x = Inches(0.5 + col * 4)
        cell = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4), Inches(0.48))
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
        cell.line.color.rgb = ACCENT_GRAY

        txt = s4.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.28)
        )
        txt.text_frame.text = val
        txt.text_frame.paragraphs[0].font.size = Pt(11)
        txt.text_frame.paragraphs[0].font.color.rgb = (
            DARK_GRAY if col == 0 else (SUCCESS_GREEN if col == 1 else MEDIUM_GRAY)
        )
        if col > 0:
            txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# USP highlight
usp_box = s4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5), Inches(12.333), Inches(0.8)
)
usp_box.fill.solid()
usp_box.fill.fore_color.rgb = RGBColor(232, 240, 254)
usp_box.line.color.rgb = PRIMARY_BLUE

usp_txt = s4.shapes.add_textbox(Inches(1), Inches(5.15), Inches(11.333), Inches(0.5))
usp_txt.text_frame.text = (
    "🎯 UNIQUE: Forgotten luggage delivery service - No competitor offers this!"
)
usp_txt.text_frame.paragraphs[0].font.size = Pt(14)
usp_txt.text_frame.paragraphs[0].font.bold = True
usp_txt.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE

add_notes(
    s4,
    "BagDrop beats traditional options on every metric. Half the price at ₹50 for 6 hours. Better security with OTP/NFC. Multiple payment options. And our unique delivery service for forgotten luggage - the key differentiator.",
)

print("Slides 1-4 created with modern design!")

# Save progress
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4_Redesigned.pptx"
)
print("Saved first 4 redesigned slides")
