#!/usr/bin/env python3
"""Complete BagDrop PPT V4 - Slides 18-20 (Final)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load existing presentation
prs = Presentation(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4.pptx"
)

YELLOW = RGBColor(255, 215, 0)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)
DARK_GRAY = RGBColor(64, 64, 64)


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK
    if subtitle:
        tf.add_paragraph().text = subtitle
        tf.paragraphs[1].font.size = Pt(16)
        tf.paragraphs[1].font.color.rgb = GRAY


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ===================== SLIDE 18: GOVERNMENT SUPPORT =====================
s18 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s18)
add_title(s18, "Government Support", "Leveraging startup policies")

schemes = [
    ("Gujarat Startup Policy", "Up to ₹30 Lakh seed funding"),
    ("Startup India", "Tax holiday for 3 years"),
    ("GST Benefits", "Input credit on tech expenses"),
    ("MUDRA/MSME Loans", "₹10-50 Lakh low-interest debt"),
    ("iCreate/AIC Support", "Free mentoring, co-working space"),
]

for i, (scheme, benefit) in enumerate(schemes):
    y = Inches(1.4 + i * 1.0)

    scheme_box = s18.shapes.add_textbox(Inches(0.5), y, Inches(4.5), Inches(0.4))
    scheme_box.text_frame.text = scheme
    scheme_box.text_frame.paragraphs[0].font.size = Pt(16)
    scheme_box.text_frame.paragraphs[0].font.bold = True
    scheme_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    arrow = s18.shapes.add_textbox(Inches(5.2), y, Inches(0.5), Inches(0.4))
    arrow.text_frame.text = "→"
    arrow.text_frame.paragraphs[0].font.size = Pt(18)
    arrow.text_frame.paragraphs[0].font.color.rgb = YELLOW

    benefit_box = s18.shapes.add_textbox(Inches(5.8), y, Inches(6.5), Inches(0.4))
    benefit_box.text_frame.text = benefit
    benefit_box.text_frame.paragraphs[0].font.size = Pt(13)
    benefit_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

# Kutch advantage
kutch_box = s18.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6), Inches(6), Inches(1.2)
)
kutch_box.fill.solid()
kutch_box.fill.fore_color.rgb = RGBColor(255, 250, 220)
kutch_box.line.color.rgb = YELLOW

kutch_txt = s18.shapes.add_textbox(Inches(0.7), Inches(6.15), Inches(5.6), Inches(0.9))
kutch_tf = kutch_txt.text_frame
kutch_tf.text = "Kutch Advantage"
kutch_tf.paragraphs[0].font.size = Pt(16)
kutch_tf.paragraphs[0].font.bold = True
kutch_tf.paragraphs[0].font.color.rgb = BLACK
kutch_tf.add_paragraph().text = "Gujarat government actively supports"
kutch_tf.add_paragraph().text = "tourism & logistics startups"
for p in kutch_tf.paragraphs[1:]:
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

add_notes(
    s18,
    "We leverage multiple government schemes. Gujarat Startup Policy offers up to ₹30 Lakh. Startup India provides 3-year tax holiday. GST input credit reduces costs. MUDRA/MSME loans available. iCreate and AIC provide free mentoring. Being in Kutch gives us access to Gujarat's active support for tourism and logistics startups.",
)

# ===================== SLIDE 19: CALL TO ACTION =====================
s19 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s19)
add_title(s19, "Join Us", "Be part of India's travel revolution")

# Investment summary
inv_box = s19.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(3)
)
inv_box.fill.solid()
inv_box.fill.fore_color.rgb = YELLOW
inv_box.line.fill.background()

inv_txt = s19.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.6), Inches(2.6))
inv_tf = inv_txt.text_frame
inv_tf.text = "Investment Opportunity"
inv_tf.paragraphs[0].font.size = Pt(20)
inv_tf.paragraphs[0].font.bold = True
inv_tf.paragraphs[0].font.color.rgb = BLACK
inv_tf.add_paragraph().text = ""
inv_tf.add_paragraph().text = "Seed Funding: ₹36 Lakh"
inv_tf.add_paragraph().text = "Expected Return: 5-10x"
inv_tf.add_paragraph().text = "Break-even: Month 10-11"
inv_tf.add_paragraph().text = "Market: ₹2,000+ Crore TAM"
for i, p in enumerate(inv_tf.paragraphs[1:], 1):
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY if i > 1 else BLACK
    if i == 2:
        p.font.bold = True

# Why invest
why_items = [
    "✓ Solves real pain point (70% travelers affected)",
    "✓ Unique USP: Forgotten luggage delivery",
    "✓ Asset-light, scalable model",
    "✓ Experienced team with local expertise",
    "✓ Government support & startup benefits",
    "✓ Clear path to profitability & exit",
]

for i, item in enumerate(why_items):
    y = Inches(1.4 + i * 0.7)
    item_box = s19.shapes.add_textbox(Inches(7), y, Inches(5.8), Inches(0.5))
    item_box.text_frame.text = item
    item_box.text_frame.paragraphs[0].font.size = Pt(13)
    item_box.text_frame.paragraphs[0].font.color.rgb = BLACK

# Contact
c_box = s19.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5), Inches(12.333), Inches(2)
)
c_box.fill.solid()
c_box.fill.fore_color.rgb = BLACK
c_box.line.fill.background()

c_txt = s19.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.333), Inches(1.6))
c_tf = c_txt.text_frame
c_tf.text = "TeamBack | Kutch, Gujarat"
c_tf.paragraphs[0].font.size = Pt(24)
c_tf.paragraphs[0].font.bold = True
c_tf.paragraphs[0].font.color.rgb = YELLOW
c_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
c_tf.add_paragraph().text = "Ready to transform travel in India"
c_tf.paragraphs[1].font.size = Pt(18)
c_tf.paragraphs[1].font.color.rgb = WHITE
c_tf.paragraphs[1].alignment = PP_ALIGN.CENTER
c_tf.add_paragraph().text = ""
c_tf.add_paragraph().text = "Thank You! Questions?"
c_tf.paragraphs[3].font.size = Pt(20)
c_tf.paragraphs[3].font.bold = True
c_tf.paragraphs[3].font.color.rgb = YELLOW
c_tf.paragraphs[3].alignment = PP_ALIGN.CENTER

add_notes(
    s19,
    "This is our call to action. We're seeking ₹36 Lakh seed funding with expected 5-10x returns. Break-even at Month 10-11. Why invest? We solve a real problem affecting 70% of travelers. We have a unique USP with forgotten luggage delivery. Asset-light scalable model. Experienced local team. Government support. Clear path to exit. Join TeamBack in Kutch, Gujarat to transform travel in India!",
)

# ===================== SLIDE 20: THANK YOU / Q&A =====================
s20 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s20)

# Large thank you
thanks_box = s20.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
thanks_tf = thanks_box.text_frame
thanks_tf.text = "Thank You!"
thanks_tf.paragraphs[0].font.size = Pt(72)
thanks_tf.paragraphs[0].font.bold = True
thanks_tf.paragraphs[0].font.color.rgb = BLACK
thanks_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = s20.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(0.8))
sub_tf = sub_box.text_frame
sub_tf.text = "BagDrop: India's Smart Luggage Storage Network"
sub_tf.paragraphs[0].font.size = Pt(28)
sub_tf.paragraphs[0].font.color.rgb = GRAY
sub_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Team info
team_box = s20.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1))
team_tf = team_box.text_frame
team_tf.text = "TeamBack | Kutch, Gujarat | January 2026"
team_tf.add_paragraph().text = "Questions & Answers"
for p in team_tf.paragraphs:
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

# Gold accent
accent = s20.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.2), Inches(5.333), Inches(0.15)
)
accent.fill.solid()
accent.fill.fore_color.rgb = YELLOW
accent.line.fill.background()

add_notes(
    s20,
    "Thank you for your attention! This concludes our presentation. We're happy to answer any questions about BagDrop, our business model, financials, or partnership opportunities. Contact TeamBack in Kutch, Gujarat. Let's revolutionize travel in India together!",
)

# Save final presentation
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4.pptx"
)
print("COMPLETE! All 20 slides created successfully!")
print("File saved: BagDrop_Presentation_V4.pptx")
