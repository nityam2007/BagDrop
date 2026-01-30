# -*- coding: utf-8 -*-
"""
BagDrop Presentation Generator - IMPROVED VERSION
Swiss Style: Bold Typography, Yellow Accent, Black/White Theme
Fonts: Arial Black (headings), Montserrat (body)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Create presentation (16:9 widescreen)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 215, 0)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_GRAY = RGBColor(120, 120, 120)

# Fonts
HEADING_FONT = "Arial Black"
BODY_FONT = "Montserrat"  # Google Font - works in Google Slides

# Helper functions
def add_rectangle(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size, font_color, bold=False, font_name=BODY_FONT, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_heading(slide, left, top, width, height, text, font_size, font_color, align=PP_ALIGN.LEFT):
    return add_text_box(slide, left, top, width, height, text, font_size, font_color, bold=True, font_name=HEADING_FONT, align=align)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ============================================================
# SLIDE 1: TITLE (GOOD - Keep as is)
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1, WHITE)

add_rectangle(slide1, Inches(0), Inches(2), Inches(2), Inches(3), YELLOW)
add_heading(slide1, Inches(1.5), Inches(2), Inches(10), Inches(1.5), 
            "BAGDROP", 96, BLACK)
add_text_box(slide1, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8), 
             "SMART LUGGAGE STORAGE & DELIVERY", 32, DARK_GRAY, bold=True)
add_text_box(slide1, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5), 
             "Secure. Contactless. Affordable.", 24, YELLOW)

# ============================================================
# SLIDE 2: PROBLEM (IMPROVED)
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, WHITE)

# Yellow accent - top right
add_rectangle(slide2, Inches(10.5), Inches(0), Inches(2.833), Inches(1.8), YELLOW)

# Heading
add_heading(slide2, Inches(0.8), Inches(0.8), Inches(9), Inches(1), 
            "THE PROBLEM", 64, BLACK)

# Subheading
add_text_box(slide2, Inches(0.8), Inches(1.8), Inches(9), Inches(0.6), 
             "Travelers face a daily struggle with luggage", 22, LIGHT_GRAY)

# Problem cards - more visual layout
# Card 1
add_rectangle(slide2, Inches(0.5), Inches(2.8), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245))
add_text_box(slide2, Inches(0.8), Inches(3.1), Inches(5.4), Inches(1), 
             "EARLY CHECKOUT, LATE DEPARTURE\nHotels force checkout but flights are hours away", 18, DARK_GRAY)

# Card 2
add_rectangle(slide2, Inches(6.8), Inches(2.8), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245))
add_text_box(slide2, Inches(7.1), Inches(3.1), Inches(5.4), Inches(1), 
             "HEAVY BAGS, RUINED PLANS\nDragging luggage kills the travel experience", 18, DARK_GRAY)

# Card 3
add_rectangle(slide2, Inches(0.5), Inches(4.5), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245))
add_text_box(slide2, Inches(0.8), Inches(4.8), Inches(5.4), Inches(1), 
             "CLOAKROOMS? OVERPRICED\nRailway lockers are expensive, crowded, or broken", 18, DARK_GRAY)

# Card 4
add_rectangle(slide2, Inches(6.8), Inches(4.5), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245))
add_text_box(slide2, Inches(7.1), Inches(4.8), Inches(5.4), Inches(1), 
             "FORGOT YOUR BAG?\nNo easy way to retrieve forgotten luggage", 18, DARK_GRAY)

# Bottom stat
add_text_box(slide2, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5), 
             "70% of travelers have faced luggage hassles during their trip", 20, YELLOW, bold=True)

# ============================================================
# SLIDE 3: SOLUTION (IMPROVED)
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BLACK)

# Yellow accent stripe
add_rectangle(slide3, Inches(0), Inches(0), Inches(0.3), Inches(7.5), YELLOW)

# Heading
add_heading(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(1), 
            "THE SOLUTION", 64, YELLOW)

# Main value prop
add_text_box(slide3, Inches(0.8), Inches(1.6), Inches(11), Inches(1), 
             "Store your bags anywhere. Pick them up anytime.", 32, WHITE, bold=True)

# Solution points with icons (text-based)
add_rectangle(slide3, Inches(0.8), Inches(3), Inches(0.15), Inches(0.8), YELLOW)
add_text_box(slide3, Inches(1.2), Inches(3), Inches(10), Inches(0.9), 
             "NEIGHBORHOOD NETWORK\nWe partner with verified local shops, cafes & hotels as storage points", 20, WHITE)

add_rectangle(slide3, Inches(0.8), Inches(4.2), Inches(0.15), Inches(0.8), YELLOW)
add_text_box(slide3, Inches(1.2), Inches(4.2), Inches(10), Inches(0.9), 
             "INSTANT ACCESS\nBook via app, drop your bag, get a secure OTP or NFC token", 20, WHITE)

add_rectangle(slide3, Inches(0.8), Inches(5.4), Inches(0.15), Inches(0.8), YELLOW)
add_text_box(slide3, Inches(1.2), Inches(5.4), Inches(10), Inches(0.9), 
             "FORGOTTEN LUGGAGE? DELIVERED\nLeft something behind? We'll courier it to you anywhere in India", 20, WHITE)

# ============================================================
# SLIDE 4: VALUE PROPOSITION (IMPROVED)
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, YELLOW)

# Heading
add_heading(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(1.2), 
            "WHY CHOOSE US?", 72, BLACK)

# Value props in a cleaner format
add_text_box(slide4, Inches(0.8), Inches(2), Inches(5.5), Inches(0.5), 
             "AFFORDABLE", 28, BLACK, bold=True)
add_text_box(slide4, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.6), 
             "Starting at just Rs.50 for 6 hours", 22, DARK_GRAY)

add_text_box(slide4, Inches(7), Inches(2), Inches(5.5), Inches(0.5), 
             "FLEXIBLE PAYMENTS", 28, BLACK, bold=True)
add_text_box(slide4, Inches(7), Inches(2.5), Inches(5.5), Inches(0.6), 
             "UPI, Cards, NFC - pay your way", 22, DARK_GRAY)

add_text_box(slide4, Inches(0.8), Inches(3.6), Inches(5.5), Inches(0.5), 
             "SECURE ACCESS", 28, BLACK, bold=True)
add_text_box(slide4, Inches(0.8), Inches(4.1), Inches(5.5), Inches(0.6), 
             "OTP & NFC tokens for safe retrieval", 22, DARK_GRAY)

add_text_box(slide4, Inches(7), Inches(3.6), Inches(5.5), Inches(0.5), 
             "DELIVERY BACKUP", 28, BLACK, bold=True)
add_text_box(slide4, Inches(7), Inches(4.1), Inches(5.5), Inches(0.6), 
             "Forgot your bag? We deliver it to you", 22, DARK_GRAY)

add_text_box(slide4, Inches(0.8), Inches(5.2), Inches(5.5), Inches(0.5), 
             "PAN-INDIA NETWORK", 28, BLACK, bold=True)
add_text_box(slide4, Inches(0.8), Inches(5.7), Inches(5.5), Inches(0.6), 
             "Growing presence across major cities", 22, DARK_GRAY)

add_text_box(slide4, Inches(7), Inches(5.2), Inches(5.5), Inches(0.5), 
             "24/7 SUPPORT", 28, BLACK, bold=True)
add_text_box(slide4, Inches(7), Inches(5.7), Inches(5.5), Inches(0.6), 
             "Help whenever you need it", 22, DARK_GRAY)

# ============================================================
# SLIDE 5: TARGET AUDIENCE (GOOD - Keep as is)
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, WHITE)

add_rectangle(slide5, Inches(0), Inches(0), Inches(0.5), Inches(7.5), YELLOW)
add_heading(slide5, Inches(1), Inches(0.5), Inches(11), Inches(1), 
            "TARGET AUDIENCE", 64, BLACK)

audience = """->  Tourists & Backpackers

->  Train & Bus Travelers

->  Students & Daily Commuters

->  Business Travelers

->  Anyone with forgotten luggage needs"""

add_text_box(slide5, Inches(1), Inches(2), Inches(11), Inches(5), 
             audience, 32, DARK_GRAY)

# ============================================================
# SLIDE 6: KEY FEATURES (GOOD - Keep as is)
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, WHITE)

add_heading(slide6, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), 
            "KEY FEATURES", 56, BLACK)

# Row 1
add_rectangle(slide6, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2), YELLOW)
add_text_box(slide6, Inches(0.7), Inches(2), Inches(3.4), Inches(1.2), 
             "APP BOOKING\n& MAP VIEW", 24, BLACK, bold=True, align=PP_ALIGN.CENTER)

add_rectangle(slide6, Inches(4.7), Inches(1.5), Inches(3.8), Inches(2), BLACK)
add_text_box(slide6, Inches(4.9), Inches(2), Inches(3.4), Inches(1.2), 
             "SELF-SERVICE\nKIOSK", 24, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rectangle(slide6, Inches(8.9), Inches(1.5), Inches(3.8), Inches(2), YELLOW)
add_text_box(slide6, Inches(9.1), Inches(2), Inches(3.4), Inches(1.2), 
             "NFC & QR\nTOKEN", 24, BLACK, bold=True, align=PP_ALIGN.CENTER)

# Row 2
add_rectangle(slide6, Inches(0.5), Inches(4), Inches(3.8), Inches(2), BLACK)
add_text_box(slide6, Inches(0.7), Inches(4.5), Inches(3.4), Inches(1.2), 
             "REAL-TIME\nNOTIFICATIONS", 24, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rectangle(slide6, Inches(4.7), Inches(4), Inches(3.8), Inches(2), YELLOW)
add_text_box(slide6, Inches(4.9), Inches(4.5), Inches(3.4), Inches(1.2), 
             "DELIVERY\nSUPPORT", 24, BLACK, bold=True, align=PP_ALIGN.CENTER)

add_rectangle(slide6, Inches(8.9), Inches(4), Inches(3.8), Inches(2), BLACK)
add_text_box(slide6, Inches(9.1), Inches(4.5), Inches(3.4), Inches(1.2), 
             "24/7\nSUPPORT", 24, WHITE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: HOW IT WORKS (IMPROVED - Horizontal Flow)
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, WHITE)

add_heading(slide7, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), 
            "HOW IT WORKS", 56, BLACK)

add_text_box(slide7, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5), 
             "5 simple steps to freedom from your luggage", 22, LIGHT_GRAY)

# Step boxes - horizontal flow
# Step 1
add_rectangle(slide7, Inches(0.3), Inches(2.2), Inches(2.4), Inches(3.5), YELLOW)
add_text_box(slide7, Inches(0.4), Inches(2.4), Inches(2.2), Inches(0.6), 
             "01", 36, BLACK, bold=True, font_name=HEADING_FONT, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(0.4), Inches(3.1), Inches(2.2), Inches(0.5), 
             "FIND", 22, BLACK, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(0.4), Inches(3.6), Inches(2.2), Inches(1.5), 
             "Open app\nLocate nearest\nBagDrop point", 16, DARK_GRAY, align=PP_ALIGN.CENTER)

# Step 2
add_rectangle(slide7, Inches(2.9), Inches(2.2), Inches(2.4), Inches(3.5), BLACK)
add_text_box(slide7, Inches(3), Inches(2.4), Inches(2.2), Inches(0.6), 
             "02", 36, YELLOW, bold=True, font_name=HEADING_FONT, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(3), Inches(3.1), Inches(2.2), Inches(0.5), 
             "BOOK", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(3), Inches(3.6), Inches(2.2), Inches(1.5), 
             "Select time slot\nPay via UPI\nor Card", 16, RGBColor(200, 200, 200), align=PP_ALIGN.CENTER)

# Step 3
add_rectangle(slide7, Inches(5.5), Inches(2.2), Inches(2.4), Inches(3.5), YELLOW)
add_text_box(slide7, Inches(5.6), Inches(2.4), Inches(2.2), Inches(0.6), 
             "03", 36, BLACK, bold=True, font_name=HEADING_FONT, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(5.6), Inches(3.1), Inches(2.2), Inches(0.5), 
             "DROP", 22, BLACK, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(5.6), Inches(3.6), Inches(2.2), Inches(1.5), 
             "Leave bags\nGet OTP or\nNFC token", 16, DARK_GRAY, align=PP_ALIGN.CENTER)

# Step 4
add_rectangle(slide7, Inches(8.1), Inches(2.2), Inches(2.4), Inches(3.5), BLACK)
add_text_box(slide7, Inches(8.2), Inches(2.4), Inches(2.2), Inches(0.6), 
             "04", 36, YELLOW, bold=True, font_name=HEADING_FONT, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(8.2), Inches(3.1), Inches(2.2), Inches(0.5), 
             "RETRIEVE", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(8.2), Inches(3.6), Inches(2.2), Inches(1.5), 
             "Show token\nCollect bags\nTravel light!", 16, RGBColor(200, 200, 200), align=PP_ALIGN.CENTER)

# Step 5
add_rectangle(slide7, Inches(10.7), Inches(2.2), Inches(2.4), Inches(3.5), YELLOW)
add_text_box(slide7, Inches(10.8), Inches(2.4), Inches(2.2), Inches(0.6), 
             "05", 36, BLACK, bold=True, font_name=HEADING_FONT, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(10.8), Inches(3.1), Inches(2.2), Inches(0.5), 
             "DELIVER", 22, BLACK, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(10.8), Inches(3.6), Inches(2.2), Inches(1.5), 
             "Forgot bag?\nWe courier it\nto you!", 16, DARK_GRAY, align=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide7, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5), 
             "Average drop time: Under 2 minutes", 20, YELLOW, bold=True)

# ============================================================
# SLIDE 8: TECH & PAYMENTS (IMPROVED)
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, BLACK)

add_heading(slide8, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), 
            "TECHNOLOGY", 64, YELLOW)

add_text_box(slide8, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5), 
             "Built for speed, security & scale", 24, WHITE)

# Three column layout
# Column 1 - Payments
add_rectangle(slide8, Inches(0.5), Inches(2.3), Inches(3.8), Inches(4), RGBColor(30, 30, 30))
add_rectangle(slide8, Inches(0.5), Inches(2.3), Inches(3.8), Inches(0.15), YELLOW)
add_text_box(slide8, Inches(0.7), Inches(2.7), Inches(3.4), Inches(0.5), 
             "PAYMENTS", 22, YELLOW, bold=True)
add_text_box(slide8, Inches(0.7), Inches(3.3), Inches(3.4), Inches(2.5), 
             "UPI (GPay, PhonePe)\n\nCredit & Debit Cards\n\nNFC Contactless\n\nWallet Integration", 18, WHITE)

# Column 2 - Security
add_rectangle(slide8, Inches(4.7), Inches(2.3), Inches(3.8), Inches(4), RGBColor(30, 30, 30))
add_rectangle(slide8, Inches(4.7), Inches(2.3), Inches(3.8), Inches(0.15), YELLOW)
add_text_box(slide8, Inches(4.9), Inches(2.7), Inches(3.4), Inches(0.5), 
             "SECURITY", 22, YELLOW, bold=True)
add_text_box(slide8, Inches(4.9), Inches(3.3), Inches(3.4), Inches(2.5), 
             "OTP Verification\n\nNFC Token Access\n\nEnd-to-End Encryption\n\nVerified Partners", 18, WHITE)

# Column 3 - Platform
add_rectangle(slide8, Inches(8.9), Inches(2.3), Inches(3.8), Inches(4), RGBColor(30, 30, 30))
add_rectangle(slide8, Inches(8.9), Inches(2.3), Inches(3.8), Inches(0.15), YELLOW)
add_text_box(slide8, Inches(9.1), Inches(2.7), Inches(3.4), Inches(0.5), 
             "PLATFORM", 22, YELLOW, bold=True)
add_text_box(slide8, Inches(9.1), Inches(3.3), Inches(3.4), Inches(2.5), 
             "iOS & Android App\n\nSelf-Service Kiosks\n\nPartner Dashboard\n\nReal-time Analytics", 18, WHITE)

# ============================================================
# SLIDE 9: TEAM STRUCTURE (GOOD - Keep as is)
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, YELLOW)

add_heading(slide9, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), 
            "TEAM STRUCTURE", 56, BLACK)

team = """CEO / FOUNDER
Overall strategy & vision

CTO
App & platform technology

OPERATIONS MANAGER
Onboarding & quality control

LOGISTICS COORDINATOR
Delivery partner management

CUSTOMER SUPPORT
Phone & app support team"""

add_text_box(slide9, Inches(0.8), Inches(1.5), Inches(11), Inches(5.5), 
             team, 24, BLACK)

# ============================================================
# SLIDE 10: GO-TO-MARKET (IMPROVED)
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, WHITE)

add_heading(slide10, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), 
            "GO-TO-MARKET", 64, BLACK)

add_text_box(slide10, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5), 
             "Multi-channel strategy to reach travelers where they are", 20, LIGHT_GRAY)

# Marketing channels - improved cards
# Card 1 - Digital
add_rectangle(slide10, Inches(0.5), Inches(2.2), Inches(3.8), Inches(2.2), YELLOW)
add_text_box(slide10, Inches(0.7), Inches(2.5), Inches(3.4), Inches(0.5), 
             "DIGITAL ADS", 24, BLACK, bold=True)
add_text_box(slide10, Inches(0.7), Inches(3.1), Inches(3.4), Inches(1.2), 
             "Instagram, Facebook & Google targeting tourists and travelers", 16, DARK_GRAY)

# Card 2 - Influencers
add_rectangle(slide10, Inches(4.7), Inches(2.2), Inches(3.8), Inches(2.2), BLACK)
add_text_box(slide10, Inches(4.9), Inches(2.5), Inches(3.4), Inches(0.5), 
             "INFLUENCERS", 24, YELLOW, bold=True)
add_text_box(slide10, Inches(4.9), Inches(3.1), Inches(3.4), Inches(1.2), 
             "Partner with travel vloggers and content creators for authentic reach", 16, WHITE)

# Card 3 - Offline
add_rectangle(slide10, Inches(8.9), Inches(2.2), Inches(3.8), Inches(2.2), YELLOW)
add_text_box(slide10, Inches(9.1), Inches(2.5), Inches(3.4), Inches(0.5), 
             "TRANSIT HUBS", 24, BLACK, bold=True)
add_text_box(slide10, Inches(9.1), Inches(3.1), Inches(3.4), Inches(1.2), 
             "Flyers and kiosks at railway stations, airports, bus stands", 16, DARK_GRAY)

# Card 4 - Referrals
add_rectangle(slide10, Inches(0.5), Inches(4.7), Inches(5.8), Inches(1.8), BLACK)
add_text_box(slide10, Inches(0.7), Inches(5), Inches(5.4), Inches(0.5), 
             "REFERRAL PROGRAM", 24, YELLOW, bold=True)
add_text_box(slide10, Inches(0.7), Inches(5.5), Inches(5.4), Inches(0.8), 
             "Rs.50 credit for every friend referred", 18, WHITE)

# Card 5 - Partners
add_rectangle(slide10, Inches(6.8), Inches(4.7), Inches(5.8), Inches(1.8), YELLOW)
add_text_box(slide10, Inches(7), Inches(5), Inches(5.4), Inches(0.5), 
             "PARTNER ONBOARDING", 24, BLACK, bold=True)
add_text_box(slide10, Inches(7), Inches(5.5), Inches(5.4), Inches(0.8), 
             "Commission-based model for shop owners", 18, DARK_GRAY)

# ============================================================
# SLIDE 11: REVENUE MODEL (GOOD - Keep as is)
# ============================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11, BLACK)

add_heading(slide11, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), 
            "REVENUE MODEL", 56, YELLOW)

revenue = """REVENUE STREAMS
- Storage fees: Rs.50 per booking
- Delivery arrangement fees
- NFC premium token fees
- Kiosk advertisement revenue"""

add_text_box(slide11, Inches(0.8), Inches(1.5), Inches(11), Inches(2), 
             revenue, 24, WHITE)

add_heading(slide11, Inches(0.8), Inches(3.8), Inches(11), Inches(1.5), 
            "Rs.7,50,000", 96, YELLOW)

add_text_box(slide11, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8), 
             "ESTIMATED MONTHLY EARNINGS", 24, WHITE)

# ============================================================
# SLIDE 12: EXIT STRATEGY (GOOD - Keep as is)
# ============================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide12, WHITE)

add_rectangle(slide12, Inches(0), Inches(5.5), Inches(5), Inches(2), YELLOW)
add_heading(slide12, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), 
            "EXIT STRATEGY", 64, BLACK)

exit_strategy = """->  Acquisition by travel/logistics company

->  Merge with major courier platform

->  Technology IP licensing/sale

->  Expansion to metro & international markets"""

add_text_box(slide12, Inches(0.8), Inches(2), Inches(11), Inches(4), 
             exit_strategy, 32, DARK_GRAY)

# ============================================================
# SLIDE 13: ROADMAP (GOOD - Keep as is)
# ============================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide13, YELLOW)

add_heading(slide13, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), 
            "FUTURE ROADMAP", 64, BLACK)

# Timeline boxes
add_rectangle(slide13, Inches(0.5), Inches(2), Inches(2.8), Inches(1.8), BLACK)
add_text_box(slide13, Inches(0.6), Inches(2.4), Inches(2.6), Inches(1.2), 
             "Q1\n50+ CITIES", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rectangle(slide13, Inches(3.6), Inches(2), Inches(2.8), Inches(1.8), BLACK)
add_text_box(slide13, Inches(3.7), Inches(2.4), Inches(2.6), Inches(1.2), 
             "Q2\nAIRPORTS", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rectangle(slide13, Inches(6.7), Inches(2), Inches(2.8), Inches(1.8), BLACK)
add_text_box(slide13, Inches(6.8), Inches(2.4), Inches(2.6), Inches(1.2), 
             "Q3\nSUBSCRIPTIONS", 20, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rectangle(slide13, Inches(9.8), Inches(2), Inches(2.8), Inches(1.8), BLACK)
add_text_box(slide13, Inches(9.9), Inches(2.4), Inches(2.6), Inches(1.2), 
             "Q4\nSMART LOCKERS", 18, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide13, Inches(0.8), Inches(4.5), Inches(11), Inches(2), 
             "VISION: Become India's #1 luggage storage network with\nsmart locker integrations across all major transit hubs.", 
             28, BLACK)

# ============================================================
# SLIDE 14: THANK YOU (GOOD - Keep as is)
# ============================================================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide14, BLACK)

add_heading(slide14, Inches(0), Inches(2), Inches(13.333), Inches(1.5), 
            "THANK YOU", 96, YELLOW, align=PP_ALIGN.CENTER)

add_text_box(slide14, Inches(0), Inches(3.8), Inches(13.333), Inches(0.8), 
             "QUESTIONS & DISCUSSION", 32, WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide14, Inches(0), Inches(5.2), Inches(13.333), Inches(1.5), 
             "[Your Email] | [Your Phone]\nwww.bagdrop.in", 24, WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# SAVE PRESENTATION
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "BagDrop_Presentation.pptx")
prs.save(output_path)

print("\n" + "="*50)
print("  SUCCESS! Presentation Created!")
print("="*50)
print("  Total Slides: 14")
print("  Fonts: Arial Black (headings), Montserrat (body)")
print("  Design: Swiss Style (Yellow/Black/White)")
print("="*50 + "\n")
