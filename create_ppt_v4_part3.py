#!/usr/bin/env python3
"""Continue BagDrop PPT V4 - Slides 11-20"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load existing presentation
prs = Presentation(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4.pptx"
)

YELLOW = RGBColor(255, 215, 0)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(245, 245, 245)


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK
    if subtitle:
        tf.add_paragraph().text = subtitle
        tf.paragraphs[1].font.size = Pt(16)
        tf.paragraphs[1].font.color.rgb = GRAY


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ===================== SLIDE 11: MARKETING STRATEGY =====================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s11)
add_title(s11, "Marketing Strategy", "Multi-channel customer acquisition")

channels = [
    ("Digital Marketing", "₹55K/month", "Instagram, Facebook, Google Ads"),
    ("Influencer Marketing", "₹30K/month", "Travel vloggers, Gujarati creators"),
    ("On-Ground (BTL)", "₹25K/month", "Railway stations, Rann Utsav kiosks"),
    ("Referral Program", "₹5K/month", "User rewards, partner incentives"),
    ("OTA Partnerships", "Commission-based", "MakeMyTrip, Goibibo integration"),
]

for i, (channel, budget, desc) in enumerate(channels):
    y = Inches(1.4 + i * 1.0)

    ch_box = s11.shapes.add_textbox(Inches(0.5), y, Inches(3.5), Inches(0.4))
    ch_box.text_frame.text = channel
    ch_box.text_frame.paragraphs[0].font.size = Pt(16)
    ch_box.text_frame.paragraphs[0].font.bold = True
    ch_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    bud_box = s11.shapes.add_textbox(Inches(4.2), y, Inches(2), Inches(0.4))
    bud_box.text_frame.text = budget
    bud_box.text_frame.paragraphs[0].font.size = Pt(14)
    bud_box.text_frame.paragraphs[0].font.bold = True
    bud_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 150, 50)

    desc_box = s11.shapes.add_textbox(Inches(6.5), y, Inches(6), Inches(0.4))
    desc_box.text_frame.text = desc
    desc_box.text_frame.paragraphs[0].font.size = Pt(12)
    desc_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

# Total budget
total_box = s11.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(6), Inches(0.8)
)
total_box.fill.solid()
total_box.fill.fore_color.rgb = YELLOW
total_box.line.fill.background()

total_txt = s11.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(5.6), Inches(0.6))
total_tf = total_txt.text_frame
total_tf.text = "Annual Marketing Budget: ₹6.45 Lakh"
total_tf.paragraphs[0].font.size = Pt(16)
total_tf.paragraphs[0].font.bold = True
total_tf.paragraphs[0].font.color.rgb = BLACK

add_notes(
    s11,
    "Our marketing strategy uses 5 channels. Digital ads at ₹55K/month targeting travelers. Influencer partnerships with Gujarati travel vloggers. On-ground presence at railway stations and Rann Utsav. Referral program with user rewards. OTA partnerships with MakeMyTrip and Goibibo. Total annual budget: ₹6.45 Lakh.",
)

# ===================== SLIDE 12: TEAM =====================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s12)
add_title(s12, "Team TeamBack", "Experienced, passionate, local")

team = [
    ("CEO / Founder", "Strategy, Vision, Partnerships", "Full-time"),
    ("CTO", "App Development, Tech Stack", "Full-time"),
    ("Operations Manager", "Partner Onboarding, Quality", "Full-time"),
    ("Logistics Coordinator", "Delivery Partners, Couriers", "Full-time"),
    ("Customer Support", "24/7 Support, Issue Resolution", "Full-time"),
]

for i, (role, resp, time) in enumerate(team):
    y = Inches(1.4 + i * 1.0)

    # Role
    role_box = s12.shapes.add_textbox(Inches(0.5), y, Inches(3), Inches(0.4))
    role_box.text_frame.text = role
    role_box.text_frame.paragraphs[0].font.size = Pt(16)
    role_box.text_frame.paragraphs[0].font.bold = True
    role_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    # Responsibility
    resp_box = s12.shapes.add_textbox(Inches(3.8), y, Inches(5.5), Inches(0.4))
    resp_box.text_frame.text = resp
    resp_box.text_frame.paragraphs[0].font.size = Pt(12)
    resp_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # Time
    time_box = s12.shapes.add_textbox(Inches(9.5), y, Inches(2.5), Inches(0.4))
    time_box.text_frame.text = time
    time_box.text_frame.paragraphs[0].font.size = Pt(12)
    time_box.text_frame.paragraphs[0].font.color.rgb = GRAY

# Hiring plan
hire_box = s12.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(6), Inches(0.8)
)
hire_box.fill.solid()
hire_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
hire_box.line.color.rgb = RGBColor(100, 150, 200)

hire_txt = s12.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(5.6), Inches(0.6))
hire_tf = hire_txt.text_frame
hire_tf.text = "Year 1 Hiring: 8 team members"
hire_tf.paragraphs[0].font.size = Pt(14)
hire_tf.paragraphs[0].font.bold = True
hire_tf.paragraphs[0].font.color.rgb = RGBColor(50, 100, 150)

add_notes(
    s12,
    "We're TeamBack from Kutch, Gujarat. Our team has 5 core roles: CEO for strategy and partnerships, CTO for app development, Operations Manager for partner onboarding, Logistics Coordinator for delivery operations, and Customer Support for 24/7 assistance. We grow to 8 members by Year 1. Local roots give us deep market understanding.",
)

# ===================== SLIDE 13: PRODUCT FEATURES =====================
s13 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s13)
add_title(s13, "Product Features", "Mobile app + Partner dashboard")

features = [
    ("Real-time Location Tracking", "Find nearest BagDrop points"),
    ("Secure OTP/NFC Tokens", "Contactless bag retrieval"),
    ("Multiple Payment Options", "UPI, Cards, NFC wallets"),
    ("Booking Management", "Modify, extend, cancel bookings"),
    ("Forgotten Luggage Delivery", "Door-to-door courier service"),
    ("24/7 Customer Support", "In-app chat, phone support"),
    ("Partner Dashboard", "For shop owners to manage storage"),
    ("Insurance Coverage", "Up to ₹25,000 per bag"),
]

col1 = features[:4]
col2 = features[4:]

for i, (feat, desc) in enumerate(col1):
    y = Inches(1.4 + i * 1.1)

    check = s13.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.35), Inches(0.35)
    )
    check.fill.solid()
    check.fill.fore_color.rgb = YELLOW
    check.line.fill.background()

    feat_box = s13.shapes.add_textbox(Inches(1), y, Inches(5), Inches(0.4))
    feat_box.text_frame.text = feat
    feat_box.text_frame.paragraphs[0].font.size = Pt(14)
    feat_box.text_frame.paragraphs[0].font.bold = True
    feat_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    desc_box = s13.shapes.add_textbox(
        Inches(1), y + Inches(0.35), Inches(5), Inches(0.4)
    )
    desc_box.text_frame.text = desc
    desc_box.text_frame.paragraphs[0].font.size = Pt(11)
    desc_box.text_frame.paragraphs[0].font.color.rgb = GRAY

for i, (feat, desc) in enumerate(col2):
    y = Inches(1.4 + i * 1.1)

    check = s13.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7), y, Inches(0.35), Inches(0.35)
    )
    check.fill.solid()
    check.fill.fore_color.rgb = YELLOW
    check.line.fill.background()

    feat_box = s13.shapes.add_textbox(Inches(7.5), y, Inches(5), Inches(0.4))
    feat_box.text_frame.text = feat
    feat_box.text_frame.paragraphs[0].font.size = Pt(14)
    feat_box.text_frame.paragraphs[0].font.bold = True
    feat_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    desc_box = s13.shapes.add_textbox(
        Inches(7.5), y + Inches(0.35), Inches(5), Inches(0.4)
    )
    desc_box.text_frame.text = desc
    desc_box.text_frame.paragraphs[0].font.size = Pt(11)
    desc_box.text_frame.paragraphs[0].font.color.rgb = GRAY

add_notes(
    s13,
    "Our app includes 8 key features. Real-time location tracking to find nearest points. Secure OTP/NFC tokens for contactless retrieval. Multiple payment options including UPI and cards. Forgotten luggage delivery - our unique feature. 24/7 support. Partner dashboard for shop owners. Insurance coverage up to ₹25,000 per bag.",
)

# ===================== SLIDE 14: PARTNER NETWORK =====================
s14 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s14)
add_title(s14, "Partner Network", "Verified local shop partners")

partner_types = [
    ("General Stores", "Near transit hubs"),
    ("Cafes & Restaurants", "Extended hours"),
    ("Budget Hotels", "Existing infrastructure"),
    ("Pharmacies", "High trust, secure"),
    ("Hostels", "Tourist-focused"),
]

for i, (ptype, loc) in enumerate(partner_types):
    y = Inches(1.4 + i * 0.8)

    type_box = s14.shapes.add_textbox(Inches(0.5), y, Inches(3.5), Inches(0.4))
    type_box.text_frame.text = ptype
    type_box.text_frame.paragraphs[0].font.size = Pt(16)
    type_box.text_frame.paragraphs[0].font.bold = True
    type_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    loc_box = s14.shapes.add_textbox(Inches(4.2), y, Inches(3), Inches(0.4))
    loc_box.text_frame.text = loc
    loc_box.text_frame.paragraphs[0].font.size = Pt(12)
    loc_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

# Partner benefits
benefits = [
    ("70% Revenue Share", "Partners keep majority"),
    ("Weekly Settlement", "Every Monday via bank"),
    ("No Lock-in Period", "Flexible partnership"),
    ("Training & Support", "Video + manual in Gujarati"),
    ("Free Marketing", "Listed on BagDrop app"),
]

for i, (benefit, desc) in enumerate(benefits):
    y = Inches(1.4 + i * 0.8)

    ben_box = s14.shapes.add_textbox(Inches(8), y, Inches(3), Inches(0.4))
    ben_box.text_frame.text = benefit
    ben_box.text_frame.paragraphs[0].font.size = Pt(14)
    ben_box.text_frame.paragraphs[0].font.bold = True
    ben_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 150, 50)

    desc_box = s14.shapes.add_textbox(Inches(11.2), y, Inches(1.8), Inches(0.4))
    desc_box.text_frame.text = desc
    desc_box.text_frame.paragraphs[0].font.size = Pt(10)
    desc_box.text_frame.paragraphs[0].font.color.rgb = GRAY

# Growth target
growth_box = s14.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(6), Inches(1.5)
)
growth_box.fill.solid()
growth_box.fill.fore_color.rgb = YELLOW
growth_box.line.fill.background()

growth_txt = s14.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(5.6), Inches(1.1))
growth_tf = growth_txt.text_frame
growth_tf.text = "Partner Growth Target"
growth_tf.paragraphs[0].font.size = Pt(18)
growth_tf.paragraphs[0].font.bold = True
growth_tf.paragraphs[0].font.color.rgb = BLACK
growth_tf.add_paragraph().text = ""
growth_tf.add_paragraph().text = "Year 1: 150 locations"
growth_tf.add_paragraph().text = "Year 3: 1,500 locations"
for p in growth_tf.paragraphs[1:]:
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY

add_notes(
    s14,
    "Our partner network includes general stores, cafes, hotels, pharmacies near transit hubs. Partners get 70% revenue share with weekly settlements. No lock-in period. Full training and support in Gujarati. We target 150 locations in Year 1, scaling to 1,500 by Year 3. This asset-light model keeps costs low.",
)

print("Slides 11-14 complete")

# ===================== SLIDE 15: EXIT STRATEGY =====================
s15 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s15)
add_title(s15, "Exit Strategy", "Multiple paths to liquidity")

exits = [
    ("Strategic Acquisition", "70% probability", "MakeMyTrip, Goibibo, Delhivery"),
    ("Merger (Logistics)", "15% probability", "Delhivery, Dunzo, Porter"),
    ("Technology Licensing", "10% probability", "White-label platform"),
    ("IPO (Long-term)", "5% probability", "Year 7+ horizon"),
]

for i, (exit_type, prob, example) in enumerate(exits):
    y = Inches(1.4 + i * 1.1)

    exit_box = s15.shapes.add_textbox(Inches(0.5), y, Inches(3.5), Inches(0.4))
    exit_box.text_frame.text = exit_type
    exit_box.text_frame.paragraphs[0].font.size = Pt(16)
    exit_box.text_frame.paragraphs[0].font.bold = True
    exit_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    prob_box = s15.shapes.add_textbox(Inches(4.2), y, Inches(2), Inches(0.4))
    prob_box.text_frame.text = prob
    prob_box.text_frame.paragraphs[0].font.size = Pt(13)
    prob_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 150, 50)

    ex_box = s15.shapes.add_textbox(Inches(6.5), y, Inches(5.5), Inches(0.4))
    ex_box.text_frame.text = example
    ex_box.text_frame.paragraphs[0].font.size = Pt(12)
    ex_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

# Valuation
val_box = s15.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(6), Inches(1.2)
)
val_box.fill.solid()
val_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
val_box.line.color.rgb = RGBColor(50, 150, 50)

val_txt = s15.shapes.add_textbox(Inches(0.7), Inches(6), Inches(5.6), Inches(0.8))
val_tf = val_txt.text_frame
val_tf.text = "Target Valuation: ₹50+ Crore"
val_tf.paragraphs[0].font.size = Pt(18)
val_tf.paragraphs[0].font.bold = True
val_tf.paragraphs[0].font.color.rgb = RGBColor(50, 150, 50)
val_tf.add_paragraph().text = "By Year 5 (3-8x ARR multiple)"
val_tf.paragraphs[1].font.size = Pt(12)
val_tf.paragraphs[1].font.color.rgb = DARK_GRAY

add_notes(
    s15,
    "We have 4 exit paths. Most likely is strategic acquisition by MakeMyTrip, Goibibo, or Delhivery at 70% probability. Merger with logistics players at 15%. Technology licensing at 10%. IPO is a 5% long-term option. Target valuation: ₹50+ Crore by Year 5 based on 3-8x ARR multiples.",
)

# ===================== SLIDE 16: RISK MITIGATION =====================
s16 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s16)
add_title(s16, "Risk Mitigation", "Addressing potential challenges")

risks = [
    ("Low Bookings", "Aggressive referral program, dynamic pricing"),
    ("Partner Defaults", "Weekly settlements, ₹5K security deposit"),
    ("Fraud/Theft", "Insurance, partner verification, CCTV"),
    ("Competition Entry", "Focus on USP, partner lock-in"),
    ("Payment Disputes", "Clear T&C, escrow for high-value items"),
]

for i, (risk, mitigation) in enumerate(risks):
    y = Inches(1.4 + i * 1.0)

    risk_box = s16.shapes.add_textbox(Inches(0.5), y, Inches(3.5), Inches(0.4))
    risk_box.text_frame.text = risk
    risk_box.text_frame.paragraphs[0].font.size = Pt(15)
    risk_box.text_frame.paragraphs[0].font.bold = True
    risk_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 50, 50)

    arrow = s16.shapes.add_textbox(Inches(4.2), y, Inches(0.5), Inches(0.4))
    arrow.text_frame.text = "→"
    arrow.text_frame.paragraphs[0].font.size = Pt(16)
    arrow.text_frame.paragraphs[0].font.color.rgb = YELLOW

    mit_box = s16.shapes.add_textbox(Inches(4.8), y, Inches(7.5), Inches(0.4))
    mit_box.text_frame.text = mitigation
    mit_box.text_frame.paragraphs[0].font.size = Pt(12)
    mit_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

# Insurance banner
ins_box = s16.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.8)
)
ins_box.fill.solid()
ins_box.fill.fore_color.rgb = BLACK
ins_box.line.fill.background()

ins_txt = s16.shapes.add_textbox(Inches(1), Inches(6.45), Inches(11.333), Inches(0.5))
ins_txt.text_frame.text = "Insurance Coverage: Public Liability ₹10L | Storage Insurance ₹25K/bag | Employee Health + Accident"
ins_txt.text_frame.paragraphs[0].font.size = Pt(13)
ins_txt.text_frame.paragraphs[0].font.color.rgb = YELLOW
ins_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_notes(
    s16,
    "We've identified key risks and mitigation strategies. Low bookings? We'll use aggressive referral programs. Partner defaults? Weekly settlements with security deposits. Theft? Insurance coverage up to ₹25K per bag plus CCTV requirements. Competition? Our delivery USP and partner lock-in. Payment disputes? Clear terms and escrow for high-value items.",
)

# ===================== SLIDE 17: MILESTONES =====================
s17 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s17)
add_title(s17, "Key Milestones", "Roadmap to success")

milestones = [
    ("Q1 2026", "App launch, 50 partners in Gujarat"),
    ("Q2 2026", "5,000 monthly bookings, expand to 3 cities"),
    ("Q3 2026", "10,000 bookings, delivery service live"),
    ("Q4 2026", "30,000 bookings, break-even achieved"),
    ("Year 2", "50 cities, 500 partners, ₹1.27 Cr revenue"),
    ("Year 3", "150 cities, 1,500 partners, ₹3.82 Cr revenue"),
    ("Year 5", "Exit readiness, ₹50+ Cr valuation"),
]

for i, (time, milestone) in enumerate(milestones):
    y = Inches(1.4 + i * 0.75)

    # Timeline dot
    dot = s17.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.05), Inches(0.25), Inches(0.25)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = YELLOW
    dot.line.fill.background()

    # Vertical line (except last)
    if i < len(milestones) - 1:
        line = s17.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.71),
            y + Inches(0.3),
            Inches(0.03),
            Inches(0.45),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GRAY
        line.line.fill.background()

    time_box = s17.shapes.add_textbox(Inches(1.2), y, Inches(2), Inches(0.4))
    time_box.text_frame.text = time
    time_box.text_frame.paragraphs[0].font.size = Pt(14)
    time_box.text_frame.paragraphs[0].font.bold = True
    time_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    mile_box = s17.shapes.add_textbox(Inches(3.5), y, Inches(9), Inches(0.4))
    mile_box.text_frame.text = milestone
    mile_box.text_frame.paragraphs[0].font.size = Pt(13)
    mile_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

add_notes(
    s17,
    "Our 5-year roadmap. Q1 2026: App launch with 50 partners. Q2: 5,000 bookings. Q3: Delivery service live. Q4: Break-even with 30,000 bookings. Year 2: 50 cities, ₹1.27 Crore revenue. Year 3: 150 cities, ₹3.82 Crore. Year 5: Exit readiness at ₹50+ Crore valuation. Clear, achievable milestones.",
)

print("Slides 15-17 complete")

# Save
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4.pptx"
)
print("Saved slides 1-17")
