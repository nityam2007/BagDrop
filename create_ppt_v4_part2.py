#!/usr/bin/env python3
"""Continue BagDrop PPT V4 - Slides 7-14"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load existing presentation
prs = Presentation(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4.pptx"
)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

YELLOW = RGBColor(255, 215, 0)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(245, 245, 245)


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK
    if subtitle:
        tf.add_paragraph().text = subtitle
        tf.paragraphs[1].font.size = Pt(16)
        tf.paragraphs[1].font.color.rgb = GRAY


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ===================== SLIDE 7: BUSINESS MODEL =====================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7)
add_title(s7, "Business Model", "Multiple revenue streams")

revenue_streams = [
    ("Storage Fees", "30% commission", "₹50-200 per booking"),
    ("Delivery Fees", "20% commission", "₹300-1,500 forgotten bags"),
    ("NFC Tokens", "100% margin", "₹199 one-time purchase"),
    ("Kiosk Ads", "100% margin", "₹5,000/month per kiosk"),
    ("Partner Subscriptions", "100% margin", "₹500/month premium"),
]

# Headers
headers = ["Revenue Stream", "Margin", "Pricing"]
for col, header in enumerate(headers):
    x = Inches(0.5 + col * 4)
    cell = s7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Inches(1.4), Inches(4), Inches(0.5)
    )
    cell.fill.solid()
    cell.fill.fore_color.rgb = YELLOW
    cell.line.fill.background()
    txt = s7.shapes.add_textbox(x + Inches(0.2), Inches(1.45), Inches(3.6), Inches(0.4))
    txt.text_frame.text = header
    txt.text_frame.paragraphs[0].font.size = Pt(14)
    txt.text_frame.paragraphs[0].font.bold = True
    txt.text_frame.paragraphs[0].font.color.rgb = BLACK

# Data rows
for row_idx, (stream, margin, price) in enumerate(revenue_streams):
    y = Inches(1.9 + row_idx * 0.6)
    data = [stream, margin, price]
    for col, val in enumerate(data):
        x = Inches(0.5 + col * 4)
        cell = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4), Inches(0.55))
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else LIGHT_GRAY
        cell.line.color.rgb = GRAY
        txt = s7.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.35)
        )
        txt.text_frame.text = val
        txt.text_frame.paragraphs[0].font.size = Pt(12)
        txt.text_frame.paragraphs[0].font.color.rgb = BLACK

# Key metric
metric_box = s7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(6), Inches(1.5)
)
metric_box.fill.solid()
metric_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
metric_box.line.color.rgb = RGBColor(50, 150, 50)

m_txt = s7.shapes.add_textbox(Inches(0.7), Inches(5.4), Inches(5.6), Inches(1.1))
m_tf = m_txt.text_frame
m_tf.text = "Pricing: ₹50 for 6 hours"
m_tf.paragraphs[0].font.size = Pt(18)
m_tf.paragraphs[0].font.bold = True
m_tf.paragraphs[0].font.color.rgb = RGBColor(50, 150, 50)
m_tf.add_paragraph().text = "Cheaper than railway cloakrooms"
m_tf.add_paragraph().text = "70% to partner, 30% to BagDrop"
for p in m_tf.paragraphs[1:]:
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

add_notes(
    s7,
    "We have 5 revenue streams. Primary is storage commission at 30%. We also take 20% on delivery fees for forgotten luggage - this is our unique offering. NFC tokens, kiosk ads, and partner subscriptions add additional revenue. Pricing starts at ₹50 for 6 hours - cheaper than railway lockers.",
)

# ===================== SLIDE 8: COMPETITIVE ADVANTAGE =====================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s8)
add_title(s8, "Competitive Advantage", "What sets us apart")

advantages = [
    (
        "Forgotten Luggage Delivery",
        "Only BagDrop offers this",
        "₹300-1,500 per delivery",
    ),
    ("OTP/NFC Security", "Contactless, secure retrieval", "No key handling"),
    ("Partner Network", "Verified local shop partners", "Lower cost than lockers"),
    ("24/7 Support", "Always available help", "Customer trust"),
    ("Local Focus", "Kutch-first approach", "Deep market understanding"),
]

for i, (adv, desc, val) in enumerate(advantages):
    y = Inches(1.4 + i * 1.0)

    check = s8.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.4), Inches(0.4)
    )
    check.fill.solid()
    check.fill.fore_color.rgb = YELLOW
    check.line.fill.background()

    check_txt = s8.shapes.add_textbox(Inches(0.55), y, Inches(0.3), Inches(0.4))
    check_txt.text_frame.text = "✓"
    check_txt.text_frame.paragraphs[0].font.size = Pt(16)
    check_txt.text_frame.paragraphs[0].font.bold = True
    check_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    adv_box = s8.shapes.add_textbox(Inches(1.1), y, Inches(4), Inches(0.4))
    adv_box.text_frame.text = adv
    adv_box.text_frame.paragraphs[0].font.size = Pt(16)
    adv_box.text_frame.paragraphs[0].font.bold = True
    adv_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    desc_box = s8.shapes.add_textbox(
        Inches(1.1), y + Inches(0.4), Inches(4), Inches(0.4)
    )
    desc_box.text_frame.text = desc
    desc_box.text_frame.paragraphs[0].font.size = Pt(11)
    desc_box.text_frame.paragraphs[0].font.color.rgb = GRAY

    val_box = s8.shapes.add_textbox(Inches(9), y, Inches(3.5), Inches(0.5))
    val_box.text_frame.text = val
    val_box.text_frame.paragraphs[0].font.size = Pt(12)
    val_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

# USP banner
usp = s8.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.8)
)
usp.fill.solid()
usp.fill.fore_color.rgb = BLACK
usp.line.fill.background()

usp_txt = s8.shapes.add_textbox(Inches(1), Inches(6.45), Inches(11.333), Inches(0.5))
usp_txt.text_frame.text = "UNIQUE SELLING POINT: Forgotten luggage delivery service - NO COMPETITOR offers this"
usp_txt.text_frame.paragraphs[0].font.size = Pt(16)
usp_txt.text_frame.paragraphs[0].font.bold = True
usp_txt.text_frame.paragraphs[0].font.color.rgb = YELLOW
usp_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_notes(
    s8,
    "What makes us different? Our forgotten luggage delivery service - NO competitor offers this. Contactless OTP/NFC security. Verified partner network instead of expensive infrastructure. 24/7 support. And our Kutch-first approach gives us deep local understanding. Our USP is clear: we're the only ones who will deliver your forgotten bags anywhere in India.",
)

# ===================== SLIDE 9: REVENUE PROJECTIONS =====================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s9)
add_title(s9, "Financial Projections", "3-year revenue forecast")

# Year 1 monthly table (simplified)
months = ["M1", "M3", "M6", "M9", "M12"]
bookings = [500, 2000, 7500, 15000, 30000]
revenues = ["₹9K", "₹39K", "₹1.57L", "₹3.6L", "₹7.65L"]

# Headers
headers = ["Month", "Bookings", "Revenue"]
for col, h in enumerate(headers):
    x = Inches(0.5 + col * 3)
    cell = s9.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Inches(1.4), Inches(3), Inches(0.4)
    )
    cell.fill.solid()
    cell.fill.fore_color.rgb = YELLOW
    cell.line.fill.background()
    txt = s9.shapes.add_textbox(
        x + Inches(0.2), Inches(1.42), Inches(2.6), Inches(0.35)
    )
    txt.text_frame.text = h
    txt.text_frame.paragraphs[0].font.size = Pt(13)
    txt.text_frame.paragraphs[0].font.bold = True

# Data
for i, (m, b, r) in enumerate(zip(months, bookings, revenues)):
    y = Inches(1.8 + i * 0.45)
    data = [m, str(b), r]
    for col, val in enumerate(data):
        x = Inches(0.5 + col * 3)
        cell = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3), Inches(0.42))
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY
        cell.line.color.rgb = GRAY
        txt = s9.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.05), Inches(2.6), Inches(0.32)
        )
        txt.text_frame.text = val
        txt.text_frame.paragraphs[0].font.size = Pt(12)
        txt.text_frame.paragraphs[0].font.color.rgb = BLACK

# 3-year projection
years_data = [
    ("Year 1", "1.32 Lakh", "₹36.8 Lakh", "8 people", "150 locations"),
    ("Year 2", "5 Lakh", "₹1.27 Crore", "25 people", "500 locations"),
    ("Year 3", "15 Lakh", "₹3.82 Crore", "60 people", "1,500 locations"),
]

y_headers = ["", "Bookings", "Revenue", "Team", "Locations"]
for col, h in enumerate(y_headers):
    x = Inches(5 + col * 1.9)
    cell = s9.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Inches(1.4), Inches(1.9), Inches(0.4)
    )
    cell.fill.solid()
    cell.fill.fore_color.rgb = YELLOW
    cell.line.fill.background()
    txt = s9.shapes.add_textbox(
        x + Inches(0.1), Inches(1.42), Inches(1.7), Inches(0.35)
    )
    txt.text_frame.text = h
    txt.text_frame.paragraphs[0].font.size = Pt(11)
    txt.text_frame.paragraphs[0].font.bold = True

for row_idx, (year, bk, rev, team, locs) in enumerate(years_data):
    y = Inches(1.8 + row_idx * 0.5)
    data = [year, bk, rev, team, locs]
    for col, val in enumerate(data):
        x = Inches(5 + col * 1.9)
        cell = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.9), Inches(0.47))
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else LIGHT_GRAY
        cell.line.color.rgb = GRAY
        txt = s9.shapes.add_textbox(
            x + Inches(0.1), y + Inches(0.08), Inches(1.7), Inches(0.3)
        )
        txt.text_frame.text = val
        txt.text_frame.paragraphs[0].font.size = Pt(10)
        txt.text_frame.paragraphs[0].font.color.rgb = BLACK

# Break-even note
be_box = s9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(4.5), Inches(1)
)
be_box.fill.solid()
be_box.fill.fore_color.rgb = RGBColor(235, 255, 235)
be_box.line.color.rgb = RGBColor(50, 150, 50)

be_txt = s9.shapes.add_textbox(Inches(0.7), Inches(4.65), Inches(4.1), Inches(0.7))
be_tf = be_txt.text_frame
be_tf.text = "Break-even: Month 10-11"
be_tf.paragraphs[0].font.size = Pt(14)
be_tf.paragraphs[0].font.bold = True
be_tf.paragraphs[0].font.color.rgb = RGBColor(50, 150, 50)
be_tf.add_paragraph().text = "Profitable from Q4 onwards"
be_tf.paragraphs[1].font.size = Pt(11)
be_tf.paragraphs[1].font.color.rgb = DARK_GRAY

add_notes(
    s9,
    "Our financial projections show clear growth. Starting with 500 bookings in Month 1, reaching 30,000 by Month 12. Year 1 revenue: ₹36.8 Lakh. By Year 3: ₹3.82 Crore with 15 Lakh bookings. We break even at Month 10-11 and become profitable from Q4. 150 locations in Year 1, scaling to 1,500 by Year 3.",
)

# ===================== SLIDE 10: FUNDING ASK =====================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s10)
add_title(s10, "The Ask", "Seed funding requirements")

# Total ask - BIG
ask_box = s10.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(1.8)
)
ask_box.fill.solid()
ask_box.fill.fore_color.rgb = YELLOW
ask_box.line.fill.background()

ask_txt = s10.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.6), Inches(1.4))
ask_tf = ask_txt.text_frame
ask_tf.text = "SEED FUNDING"
ask_tf.paragraphs[0].font.size = Pt(18)
ask_tf.paragraphs[0].font.bold = True
ask_tf.paragraphs[0].font.color.rgb = BLACK
ask_tf.add_paragraph().text = ""
ask_tf.add_paragraph().text = "₹36 Lakh"
ask_tf.paragraphs[2].font.size = Pt(48)
ask_tf.paragraphs[2].font.bold = True
ask_tf.paragraphs[2].font.color.rgb = BLACK
ask_tf.paragraphs[2].alignment = PP_ALIGN.CENTER
ask_tf.add_paragraph().text = "(~$430,000 USD)"
ask_tf.paragraphs[3].font.size = Pt(14)
ask_tf.paragraphs[3].font.color.rgb = DARK_GRAY
ask_tf.paragraphs[3].alignment = PP_ALIGN.CENTER

# Breakdown
breakdown = [
    ("Tech Development", "₹5 Lakh", "App, dashboard, payment gateway"),
    ("6-Month Runway", "₹18 Lakh", "Salaries + operations"),
    ("Marketing", "₹6.5 Lakh", "Digital + on-ground campaigns"),
    ("Equipment & Setup", "₹3 Lakh", "Laptops, kiosks, office"),
    ("Contingency", "₹3.25 Lakh", "10% buffer"),
]

for i, (item, amount, desc) in enumerate(breakdown):
    y = Inches(1.4 + i * 0.9)

    item_box = s10.shapes.add_textbox(Inches(7), y, Inches(3), Inches(0.4))
    item_box.text_frame.text = item
    item_box.text_frame.paragraphs[0].font.size = Pt(14)
    item_box.text_frame.paragraphs[0].font.bold = True
    item_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    amt_box = s10.shapes.add_textbox(Inches(10.2), y, Inches(2.5), Inches(0.4))
    amt_box.text_frame.text = amount
    amt_box.text_frame.paragraphs[0].font.size = Pt(14)
    amt_box.text_frame.paragraphs[0].font.bold = True
    amt_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 150, 50)

    desc_box = s10.shapes.add_textbox(
        Inches(7), y + Inches(0.35), Inches(5.5), Inches(0.4)
    )
    desc_box.text_frame.text = desc
    desc_box.text_frame.paragraphs[0].font.size = Pt(11)
    desc_box.text_frame.paragraphs[0].font.color.rgb = GRAY

# ROI
roi_box = s10.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(6), Inches(1.5)
)
roi_box.fill.solid()
roi_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
roi_box.line.color.rgb = RGBColor(100, 150, 200)

roi_txt = s10.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(5.6), Inches(1.1))
roi_tf = roi_txt.text_frame
roi_tf.text = "Expected Returns"
roi_tf.paragraphs[0].font.size = Pt(16)
roi_tf.paragraphs[0].font.bold = True
roi_tf.paragraphs[0].font.color.rgb = RGBColor(50, 100, 150)
roi_tf.add_paragraph().text = ""
roi_tf.add_paragraph().text = "5-10x return for seed investors"
roi_tf.add_paragraph().text = "By Year 5: ₹50+ Crore valuation"
for p in roi_tf.paragraphs[1:]:
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

add_notes(
    s10,
    "We're seeking ₹36 Lakh in seed funding. This breaks down to: ₹5L for tech development, ₹18L for 6-month runway, ₹6.5L for marketing to acquire users, ₹3L for equipment, and 10% contingency. Expected returns are 5-10x for seed investors by Year 5 with a target valuation of ₹50+ Crore.",
)

print("Slides 7-10 complete")

# Save
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4.pptx"
)
print("Saved slides 1-10")
