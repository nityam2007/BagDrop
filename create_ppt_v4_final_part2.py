#!/usr/bin/env python3
"""Continue BagDrop PPT V4 Final - Slides 8-14"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Load existing presentation
prs = Presentation(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4_Final.pptx"
)

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 215, 0)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_GRAY = RGBColor(120, 120, 120)
GREEN = RGBColor(52, 168, 83)

HEADING_FONT = "Arial Black"
BODY_FONT = "Montserrat"


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_heading(
    slide, left, top, width, height, text, font_size, color, align=PP_ALIGN.LEFT
):
    return add_text(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size,
        color,
        bold=True,
        font=HEADING_FONT,
        align=align,
    )


def set_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================
# SLIDE 8: HOW IT WORKS
# ============================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8, WHITE)
add_heading(
    s8, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), "HOW IT WORKS", 56, BLACK
)
add_text(
    s8,
    Inches(0.8),
    Inches(1.1),
    Inches(11),
    Inches(0.5),
    "5 simple steps to freedom from your luggage",
    22,
    LIGHT_GRAY,
)

# Step boxes
steps = [
    ("01", "FIND", "Open app\nLocate nearest\nBagDrop point"),
    ("02", "BOOK", "Select time slot\nPay via UPI\nor Card"),
    ("03", "DROP", "Leave bags\nGet OTP or\nNFC token"),
    ("04", "RETRIEVE", "Show token\nCollect bags\nTravel light!"),
    ("05", "DELIVER", "Forgot bag?\nWe courier it\nto you!"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.3 + i * 2.6)
    color = YELLOW if i % 2 == 0 else BLACK
    text_color = BLACK if i % 2 == 0 else YELLOW
    desc_color = DARK_GRAY if i % 2 == 0 else RGBColor(200, 200, 200)

    add_rect(s8, x, Inches(2.2), Inches(2.4), Inches(3.5), color)
    add_text(
        s8,
        x,
        Inches(2.4),
        Inches(2.4),
        Inches(0.6),
        num,
        36,
        text_color,
        bold=True,
        font=HEADING_FONT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s8,
        x,
        Inches(3.1),
        Inches(2.4),
        Inches(0.5),
        title,
        22,
        text_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s8,
        x,
        Inches(3.6),
        Inches(2.4),
        Inches(1.5),
        desc,
        16,
        desc_color,
        align=PP_ALIGN.CENTER,
    )

add_text(
    s8,
    Inches(0.8),
    Inches(6.2),
    Inches(11),
    Inches(0.5),
    "Average drop time: Under 2 minutes",
    20,
    YELLOW,
    bold=True,
)
add_notes(
    s8,
    "Simple 5-step process. Find nearby location. Book and pay. Drop bags with OTP/NFC token. Retrieve easily. And unique delivery service for forgotten luggage.",
)

# ============================================
# SLIDE 9: TECHNOLOGY
# ============================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9, BLACK)
add_heading(
    s9, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "TECHNOLOGY", 64, YELLOW
)
add_text(
    s9,
    Inches(0.8),
    Inches(1.4),
    Inches(11),
    Inches(0.5),
    "Built for speed, security & scale",
    24,
    WHITE,
)

# Three columns
columns = [
    (
        "PAYMENTS",
        "UPI (GPay, PhonePe)\n\nCredit & Debit Cards\n\nNFC Contactless\n\nWallet Integration",
    ),
    (
        "SECURITY",
        "OTP Verification\n\nNFC Token Access\n\nEnd-to-End Encryption\n\nVerified Partners",
    ),
    (
        "PLATFORM",
        "iOS & Android App\n\nSelf-Service Kiosks\n\nPartner Dashboard\n\nReal-time Analytics",
    ),
]

for i, (title, content) in enumerate(columns):
    x = Inches(0.5 + i * 4.2)
    add_rect(s9, x, Inches(2.3), Inches(3.8), Inches(4), RGBColor(30, 30, 30))
    add_rect(s9, x, Inches(2.3), Inches(3.8), Inches(0.15), YELLOW)
    add_text(
        s9,
        x + Inches(0.2),
        Inches(2.7),
        Inches(3.4),
        Inches(0.5),
        title,
        22,
        YELLOW,
        bold=True,
    )
    add_text(
        s9, x + Inches(0.2), Inches(3.3), Inches(3.4), Inches(2.5), content, 18, WHITE
    )

add_notes(
    s9,
    "Tech stack includes UPI, cards, NFC payments. OTP and NFC token security. iOS/Android app, self-service kiosks, partner dashboard with real-time analytics.",
)

# ============================================
# SLIDE 10: COMPETITIVE ADVANTAGE
# ============================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10, YELLOW)
add_heading(
    s10, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "WHY BAGDROP WINS", 64, BLACK
)

# Comparison table
add_rect(s10, Inches(0.5), Inches(1.7), Inches(4), Inches(0.4), BLACK)
add_text(
    s10,
    Inches(0.5),
    Inches(1.75),
    Inches(4),
    Inches(0.35),
    "FEATURE",
    13,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)

add_rect(s10, Inches(4.5), Inches(1.7), Inches(4), Inches(0.4), BLACK)
add_text(
    s10,
    Inches(4.5),
    Inches(1.75),
    Inches(4),
    Inches(0.35),
    "BAGDROP",
    13,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)

add_rect(s10, Inches(8.5), Inches(1.7), Inches(4), Inches(0.4), RGBColor(100, 100, 100))
add_text(
    s10,
    Inches(8.5),
    Inches(1.75),
    Inches(4),
    Inches(0.35),
    "OTHERS",
    13,
    WHITE,
    bold=True,
    align=PP_ALIGN.CENTER,
)

comparisons = [
    ("Pricing", "₹50/6hrs", "₹100+/day"),
    ("Security", "OTP/NFC", "Keys/Tags"),
    ("Payments", "UPI/Card/NFC", "Cash only"),
    ("Delivery", "✓ YES", "✗ NO"),
    ("Booking", "App-based", "In-person"),
    ("Support", "24/7", "Limited"),
]

for row_idx, (feature, bagdrop, others) in enumerate(comparisons):
    y = Inches(2.1 + row_idx * 0.5)
    bg = RGBColor(245, 245, 245) if row_idx % 2 == 0 else WHITE

    add_rect(s10, Inches(0.5), y, Inches(4), Inches(0.48), bg)
    add_text(
        s10,
        Inches(0.7),
        y + Inches(0.1),
        Inches(3.6),
        Inches(0.28),
        feature,
        12,
        DARK_GRAY,
    )

    add_rect(s10, Inches(4.5), y, Inches(4), Inches(0.48), bg)
    add_text(
        s10,
        Inches(4.5),
        y + Inches(0.1),
        Inches(4),
        Inches(0.28),
        bagdrop,
        12,
        GREEN if row_idx == 0 or row_idx == 3 else BLACK,
        align=PP_ALIGN.CENTER,
    )

    add_rect(s10, Inches(8.5), y, Inches(4), Inches(0.48), bg)
    add_text(
        s10,
        Inches(8.5),
        y + Inches(0.1),
        Inches(4),
        Inches(0.28),
        others,
        12,
        LIGHT_GRAY,
        align=PP_ALIGN.CENTER,
    )

# USP
add_rect(s10, Inches(0.5), Inches(5.5), Inches(12.333), Inches(1), BLACK)
add_text(
    s10,
    Inches(1),
    Inches(5.7),
    Inches(11.333),
    Inches(0.6),
    "UNIQUE: Forgotten luggage delivery - NO COMPETITOR offers this!",
    18,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)

add_notes(
    s10,
    "BagDrop beats competitors on every metric. Half the price. Better security. Multiple payment options. 24/7 support. And our unique USP: forgotten luggage delivery service.",
)

# ============================================
# SLIDE 11: FINANCIAL PROJECTIONS
# ============================================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s11, WHITE)
add_heading(
    s11,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.8),
    "FINANCIAL PROJECTIONS",
    56,
    BLACK,
)

# 3-year table
years = [
    ("YEAR 1", "1.32L bookings", "₹36.8L revenue", "8 people", "150 locations"),
    ("YEAR 2", "5L bookings", "₹1.27Cr revenue", "25 people", "500 locations"),
    ("YEAR 3", "15L bookings", "₹3.82Cr revenue", "60 people", "1,500 locations"),
]

# Headers
headers = ["", "BOOKINGS", "REVENUE", "TEAM", "LOCATIONS"]
for col, h in enumerate(headers):
    x = Inches(0.5 + col * 2.5)
    add_rect(s11, x, Inches(1.5), Inches(2.5), Inches(0.4), BLACK)
    add_text(
        s11,
        x,
        Inches(1.55),
        Inches(2.5),
        Inches(0.35),
        h,
        12,
        YELLOW,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

# Data
for row_idx, (year, bk, rev, team, locs) in enumerate(years):
    y = Inches(1.9 + row_idx * 0.6)
    data = [year, bk, rev, team, locs]
    for col, val in enumerate(data):
        x = Inches(0.5 + col * 2.5)
        bg = (
            YELLOW
            if row_idx == 2
            else (WHITE if row_idx % 2 == 0 else RGBColor(245, 245, 245))
        )
        txt_color = BLACK if row_idx == 2 else DARK_GRAY
        add_rect(s11, x, y, Inches(2.5), Inches(0.55), bg)
        add_text(
            s11,
            x + Inches(0.1),
            y + Inches(0.15),
            Inches(2.3),
            Inches(0.3),
            val,
            11,
            txt_color,
            align=PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT,
        )

# Month-by-month
add_heading(
    s11,
    Inches(0.8),
    Inches(4),
    Inches(11),
    Inches(0.6),
    "YEAR 1 MONTHLY GROWTH",
    24,
    BLACK,
)

months = [
    ("M1", "500", "₹9K"),
    ("M3", "2,000", "₹39K"),
    ("M6", "7,500", "₹1.57L"),
    ("M9", "15,000", "₹3.6L"),
    ("M12", "30,000", "₹7.65L"),
]

for i, (m, b, r) in enumerate(months):
    x = Inches(0.5 + i * 2.5)
    add_rect(s11, x, Inches(4.7), Inches(2.3), Inches(1.8), BLACK if i == 4 else YELLOW)
    txt_color = YELLOW if i == 4 else BLACK
    add_text(
        s11,
        x,
        Inches(4.9),
        Inches(2.3),
        Inches(0.5),
        m,
        28,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s11,
        x,
        Inches(5.5),
        Inches(2.3),
        Inches(0.4),
        b,
        14,
        txt_color,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s11,
        x,
        Inches(5.9),
        Inches(2.3),
        Inches(0.4),
        r,
        12,
        txt_color,
        align=PP_ALIGN.CENTER,
    )

# Break-even
add_rect(s11, Inches(0.5), Inches(6.7), Inches(5), Inches(0.5), GREEN)
add_text(
    s11,
    Inches(0.7),
    Inches(6.75),
    Inches(4.6),
    Inches(0.4),
    "Break-even: Month 10-11 | Profitable from Q4",
    14,
    WHITE,
    bold=True,
)

add_notes(
    s11,
    "Clear growth trajectory. Year 1: ₹36.8L revenue with 1.32L bookings. Year 3: ₹3.82Cr with 15L bookings. Break-even at Month 10-11. Profitable from Q4 onwards.",
)

# ============================================
# SLIDE 12: THE ASK
# ============================================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s12, BLACK)
add_heading(
    s12, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "THE ASK", 64, YELLOW
)

# Big number
add_heading(
    s12, Inches(0.8), Inches(1.5), Inches(6), Inches(1.5), "₹36 LAKH", 72, YELLOW
)
add_text(s12, Inches(0.8), Inches(3), Inches(6), Inches(0.5), "SEED FUNDING", 24, WHITE)
add_text(
    s12,
    Inches(0.8),
    Inches(3.5),
    Inches(6),
    Inches(0.4),
    "(~$430,000 USD)",
    16,
    LIGHT_GRAY,
)

# Breakdown
breakdown = [
    ("TECH DEVELOPMENT", "₹5L", "App, dashboard, payments"),
    ("6-MONTH RUNWAY", "₹18L", "Salaries + operations"),
    ("MARKETING", "₹6.5L", "Digital + on-ground"),
    ("EQUIPMENT", "₹3L", "Laptops, kiosks, office"),
    ("CONTINGENCY", "₹3.25L", "10% buffer"),
]

for i, (item, amt, desc) in enumerate(breakdown):
    y = Inches(1.5 + i * 0.9)
    add_text(s12, Inches(7), y, Inches(3.5), Inches(0.4), item, 14, WHITE, bold=True)
    add_text(s12, Inches(10.8), y, Inches(2), Inches(0.4), amt, 16, YELLOW, bold=True)
    add_text(
        s12, Inches(7), y + Inches(0.4), Inches(5.8), Inches(0.3), desc, 11, LIGHT_GRAY
    )

# ROI
add_rect(s12, Inches(0.5), Inches(5), Inches(6), Inches(1.8), YELLOW)
add_text(
    s12,
    Inches(0.7),
    Inches(5.2),
    Inches(5.6),
    Inches(0.6),
    "INVESTOR RETURNS",
    20,
    BLACK,
    bold=True,
)
add_text(
    s12,
    Inches(0.7),
    Inches(5.9),
    Inches(5.6),
    Inches(0.5),
    "5-10x Expected Return",
    24,
    BLACK,
    bold=True,
)
add_text(
    s12,
    Inches(0.7),
    Inches(6.5),
    Inches(5.6),
    Inches(0.3),
    "₹50+ Cr valuation by Year 5",
    14,
    DARK_GRAY,
)

add_notes(
    s12,
    "Seeking ₹36 Lakh seed funding. Breakdown: ₹5L tech, ₹18L runway, ₹6.5L marketing, ₹3L equipment, ₹3.25L contingency. Expected 5-10x returns with ₹50+ Cr valuation by Year 5.",
)

# ============================================
# SLIDE 13: TEAM
# ============================================
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s13, YELLOW)
add_heading(s13, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), "TEAM", 56, BLACK)

team_txt = """CEO / FOUNDER
Overall strategy & investor relations

CTO
App development & tech infrastructure

OPERATIONS MANAGER
Partner onboarding & quality control

LOGISTICS COORDINATOR
Delivery partner & courier management

CUSTOMER SUPPORT
24/7 phone & app support"""

add_text(s13, Inches(0.8), Inches(1.4), Inches(6), Inches(5.5), team_txt, 22, DARK_GRAY)

# Stats
add_rect(s13, Inches(7.5), Inches(2), Inches(5), Inches(1.2), BLACK)
add_text(
    s13,
    Inches(7.7),
    Inches(2.3),
    Inches(4.6),
    Inches(0.6),
    "8",
    48,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_text(
    s13,
    Inches(7.7),
    Inches(3),
    Inches(4.6),
    Inches(0.3),
    "Team members by Year 1",
    14,
    WHITE,
    align=PP_ALIGN.CENTER,
)

add_rect(s13, Inches(7.5), Inches(3.8), Inches(5), Inches(1.2), BLACK)
add_text(
    s13,
    Inches(7.7),
    Inches(4.1),
    Inches(4.6),
    Inches(0.6),
    "Kutch",
    36,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_text(
    s13,
    Inches(7.7),
    Inches(4.8),
    Inches(4.6),
    Inches(0.3),
    "Gujarat, India HQ",
    14,
    WHITE,
    align=PP_ALIGN.CENTER,
)

add_notes(
    s13,
    "TeamBack from Kutch, Gujarat. 5 core roles: CEO, CTO, Operations, Logistics, Support. Growing to 8 members by Year 1. Local expertise gives us deep market understanding.",
)

# ============================================
# SLIDE 14: MARKETING STRATEGY
# ============================================
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s14, WHITE)
add_heading(
    s14, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "GO-TO-MARKET", 64, BLACK
)
add_text(
    s14,
    Inches(0.8),
    Inches(1.3),
    Inches(11),
    Inches(0.5),
    "Multi-channel strategy to reach travelers",
    20,
    LIGHT_GRAY,
)

# Channel cards
channels = [
    ("DIGITAL ADS", "₹55K/month", "Instagram, Facebook, Google targeting tourists"),
    ("INFLUENCERS", "₹30K/month", "Travel vloggers & Gujarati creators"),
    ("TRANSIT HUBS", "₹25K/month", "Flyers & kiosks at stations, airports"),
]

for i, (title, budget, desc) in enumerate(channels):
    x = Inches(0.5 + i * 4.2)
    color = YELLOW if i != 1 else BLACK
    txt_color = BLACK if i != 1 else YELLOW
    desc_color = DARK_GRAY if i != 1 else RGBColor(200, 200, 200)

    add_rect(s14, x, Inches(2.2), Inches(3.9), Inches(2.2), color)
    add_text(
        s14,
        x + Inches(0.2),
        Inches(2.5),
        Inches(3.5),
        Inches(0.5),
        title,
        24,
        txt_color,
        bold=True,
    )
    add_text(
        s14,
        x + Inches(0.2),
        Inches(3.1),
        Inches(3.5),
        Inches(1.2),
        desc,
        16,
        desc_color,
    )
    add_text(
        s14,
        x + Inches(0.2),
        Inches(4.2),
        Inches(3.5),
        Inches(0.3),
        budget,
        14,
        GREEN if i != 1 else YELLOW,
        bold=True,
    )

# Bottom cards
add_rect(s14, Inches(0.5), Inches(4.7), Inches(5.8), Inches(1.8), BLACK)
add_text(
    s14,
    Inches(0.7),
    Inches(5),
    Inches(5.4),
    Inches(0.5),
    "REFERRAL PROGRAM",
    24,
    YELLOW,
    bold=True,
)
add_text(
    s14,
    Inches(0.7),
    Inches(5.5),
    Inches(5.4),
    Inches(0.8),
    "₹50 credit for every friend who signs up",
    18,
    WHITE,
)

add_rect(s14, Inches(6.8), Inches(4.7), Inches(5.8), Inches(1.8), YELLOW)
add_text(
    s14,
    Inches(7),
    Inches(5),
    Inches(5.4),
    Inches(0.5),
    "PARTNER ONBOARDING",
    24,
    BLACK,
    bold=True,
)
add_text(
    s14,
    Inches(7),
    Inches(5.5),
    Inches(5.4),
    Inches(0.8),
    "70% revenue share to shop owners",
    18,
    DARK_GRAY,
)

# Total budget
add_rect(s14, Inches(0.5), Inches(6.7), Inches(6), Inches(0.5), BLACK)
add_text(
    s14,
    Inches(0.7),
    Inches(6.75),
    Inches(5.6),
    Inches(0.4),
    "Annual Marketing Budget: ₹6.45 Lakh",
    16,
    YELLOW,
    bold=True,
)

add_notes(
    s14,
    "Multi-channel marketing: Digital ads (₹55K/month), influencers (₹30K), transit hubs (₹25K). Referral program with ₹50 credits. Partner onboarding with 70% revenue share. Total annual budget: ₹6.45 Lakh.",
)

print("Slides 8-14 complete")

# Save
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4_Final.pptx"
)
print("Saved slides 1-14")
