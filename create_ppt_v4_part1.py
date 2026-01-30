#!/usr/bin/env python3
"""
BagDrop Presentation V4 - Professional 20-Slide Deck
With Speaker Notes, Tables, and Graphics
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import io

# Create presentation with widescreen aspect ratio (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand Colors
BRAND_YELLOW = RgbColor(255, 215, 0)
BRAND_BLACK = RgbColor(0, 0, 0)
BRAND_WHITE = RgbColor(255, 255, 255)
BRAND_GRAY = RgbColor(128, 128, 128)
BRAND_DARK_GRAY = RgbColor(64, 64, 64)


def add_background(slide, color=BRAND_WHITE):
    """Add colored background to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()


def add_title_box(
    slide,
    title_text,
    subtitle_text="",
    left=Inches(0.5),
    top=Inches(0.3),
    width=Inches(12.333),
):
    """Add title and subtitle to slide"""
    # Title
    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_BLACK
    title_para.alignment = PP_ALIGN.LEFT

    # Subtitle
    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(
            left, top + Inches(0.7), width, Inches(0.4)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(16)
        subtitle_para.font.color.rgb = BRAND_GRAY
        subtitle_para.alignment = PP_ALIGN.LEFT


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to slide"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def add_table_slide(
    slide, headers, data, left=Inches(0.5), top=Inches(1.5), width=Inches(12.333)
):
    """Add a formatted table to slide"""
    rows = len(data) + 1
    cols = len(headers)

    table = slide.shapes.add_table(
        rows, cols, left, top, width, Inches(0.6 * rows)
    ).table

    # Set header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_YELLOW
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = BRAND_BLACK
        paragraph.alignment = PP_ALIGN.CENTER

    # Set data rows
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = BRAND_BLACK
            paragraph.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

            # Alternate row colors
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(245, 245, 245)

    return table


# ==========================================
# SLIDE 1: Title Slide
# ==========================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_background(slide1)

# Large title
box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.2))
frame = box.text_frame
frame.text = "BagDrop"
p = frame.paragraphs[0]
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = BRAND_BLACK
p.alignment = PP_ALIGN.CENTER

# Tagline
box2 = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.6))
frame2 = box2.text_frame
frame2.text = "India's Smart Luggage Storage Network"
p2 = frame2.paragraphs[0]
p2.font.size = Pt(24)
p2.font.color.rgb = BRAND_GRAY
p2.alignment = PP_ALIGN.CENTER

# Team info
box3 = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.8))
frame3 = box3.text_frame
frame3.text = "Team: TeamBack | Location: Kutch, Gujarat"
frame3.add_paragraph().text = "January 2026"
for p in frame3.paragraphs:
    p.font.size = Pt(14)
    p.font.color.rgb = BRAND_DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

# Gold accent bar
accent = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.5), Inches(5.333), Inches(0.1)
)
accent.fill.solid()
accent.fill.fore_color.rgb = BRAND_YELLOW
accent.line.fill.background()

add_speaker_notes(
    slide1,
    "Welcome judges and audience. Today we present BagDrop - a solution to a problem faced by 70% of travelers. We're TeamBack from Kutch, Gujarat, and we're building India's leading luggage storage network. Start with energy and confidence. This is a $15 Lakh Crore market opportunity.",
)

# ==========================================
# SLIDE 2: The Problem
# ==========================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2)
add_title_box(
    slide2, "The Traveler's Dilemma", "70% of travelers face luggage problems"
)

# Problem illustration box
prob_box = slide2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(3)
)
prob_box.fill.solid()
prob_box.fill.fore_color.rgb = RgbColor(255, 235, 235)
prob_box.line.color.rgb = RgbColor(255, 100, 100)

prob_text = slide2.shapes.add_textbox(
    Inches(0.8), Inches(1.8), Inches(5.4), Inches(2.4)
)
prob_frame = prob_text.text_frame
prob_frame.text = '"Checked out at 10 AM, but train is at 7 PM..."'
prob_frame.add_paragraph().text = ""
prob_frame.add_paragraph().text = "• Heavy bags ruin sightseeing plans"
prob_frame.add_paragraph().text = "• Railway lockers cost ₹100+/day"
prob_frame.add_paragraph().text = "• No solution for forgotten luggage"

for i, p in enumerate(prob_frame.paragraphs):
    if i == 0:
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RgbColor(180, 0, 0)
    else:
        p.font.size = Pt(14)
        p.font.color.rgb = BRAND_DARK_GRAY

# Key stats
stats_data = [
    ["5.5 Crore", "Annual domestic tourists in Gujarat"],
    ["2,400 Crore", "Annual railway passengers India"],
    ["70%", "Travelers facing luggage issues"],
]

for i, (stat, desc) in enumerate(stats_data):
    y_pos = Inches(1.5 + i * 0.9)

    stat_box = slide2.shapes.add_textbox(Inches(7), y_pos, Inches(2), Inches(0.6))
    stat_frame = stat_box.text_frame
    stat_frame.text = stat
    stat_frame.paragraphs[0].font.size = Pt(24)
    stat_frame.paragraphs[0].font.bold = True
    stat_frame.paragraphs[0].font.color.rgb = BRAND_YELLOW

    desc_box = slide2.shapes.add_textbox(
        Inches(9.2), y_pos + Inches(0.1), Inches(3.5), Inches(0.5)
    )
    desc_frame = desc_box.text_frame
    desc_frame.text = desc
    desc_frame.paragraphs[0].font.size = Pt(12)
    desc_frame.paragraphs[0].font.color.rgb = BRAND_DARK_GRAY

add_speaker_notes(
    slide2,
    "Paint the picture: Everyone has faced this. Early checkout, late departure. You drag heavy bags through cafes, temples, markets. Railway lockers are expensive, full, or broken. 5.5 Crore tourists visit Gujarat annually - all facing this problem. This is a massive pain point with no good solution.",
)

# ==========================================
# SLIDE 3: Our Solution - BagDrop
# ==========================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3)
add_title_box(slide3, "BagDrop: The Solution", "Store. Explore. Relax.")

# How it works - 5 steps
steps = [
    ("1", "FIND", "Open app, locate nearby BagDrop"),
    ("2", "BOOK", "Select slot, pay via UPI/Card"),
    ("3", "DROP", "Leave bags, get OTP/NFC token"),
    ("4", "RETRIEVE", "Show token, collect bags"),
    ("5", "DELIVER", "Forgot bag? We courier it"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(1.5 if i < 3 else 4)

    # Number circle
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BRAND_YELLOW
    circle.line.fill.background()

    num_box = slide3.shapes.add_textbox(x, y + Inches(0.1), Inches(0.6), Inches(0.4))
    num_frame = num_box.text_frame
    num_frame.text = num
    num_frame.paragraphs[0].font.size = Pt(20)
    num_frame.paragraphs[0].font.bold = True
    num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title and desc
    title_box = slide3.shapes.add_textbox(x + Inches(0.8), y, Inches(3.2), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(16)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BRAND_BLACK

    desc_box = slide3.shapes.add_textbox(
        x + Inches(0.8), y + Inches(0.4), Inches(3.2), Inches(0.8)
    )
    desc_frame = desc_box.text_frame
    desc_frame.text = desc
    desc_frame.paragraphs[0].font.size = Pt(12)
    desc_frame.paragraphs[0].font.color.rgb = BRAND_GRAY

# USP banner
usp_box = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6), Inches(12.333), Inches(1)
)
usp_box.fill.solid()
usp_box.fill.fore_color.rgb = BRAND_BLACK
usp_box.line.fill.background()

usp_text = slide3.shapes.add_textbox(
    Inches(1), Inches(6.2), Inches(11.333), Inches(0.6)
)
usp_frame = usp_text.text_frame
usp_frame.text = "Under 2 minutes average drop time | Starting at just ₹50 for 6 hours | 24/7 Support"
usp_frame.paragraphs[0].font.size = Pt(18)
usp_frame.paragraphs[0].font.color.rgb = BRAND_YELLOW
usp_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_speaker_notes(
    slide3,
    "Here's how BagDrop works - simple 5-step process. The key differentiator is our delivery service for forgotten luggage. No competitor offers this. Average drop time under 2 minutes. Price starts at just ₹50 for 6 hours - cheaper than railway cloakrooms. Emphasize the speed, affordability, and unique delivery feature.",
)

# ==========================================
# SLIDE 4: Why Kutch, Gujarat?
# ==========================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4)
add_title_box(slide4, "Why Kutch, Gujarat?", "Strategic launch location")

# Add table for location advantages
headers = ["Factor", "Advantage"]
data = [
    ["Tourist Hub", "Rann Utsav: 5 Lakh+ visitors annually"],
    ["Transit Points", "Bhuj Railway, Gandhidham Junction"],
    ["Growth", "Gujarat tourism growing at 12% CAGR"],
    ["Strategic", "Gateway to Rann, near Mundra Port"],
    ["Culture", "Handicrafts, Dholavira, year-round tourism"],
]

add_table_slide(slide4, headers, data, top=Inches(1.3))

# Seasonal highlight
season_box = slide4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5), Inches(6), Inches(2)
)
season_box.fill.solid()
season_box.fill.fore_color.rgb = RgbColor(255, 250, 220)
season_box.line.color.rgb = BRAND_YELLOW

season_text = slide4.shapes.add_textbox(
    Inches(0.8), Inches(5.2), Inches(5.4), Inches(1.6)
)
season_frame = season_text.text_frame
season_frame.text = "Rann Utsav Season (Nov-Feb)"
season_frame.add_paragraph().text = "Peak tourism with 5-7 Lakh visitors"
season_frame.add_paragraph().text = "Perfect timing for launch and validation"

season_frame.paragraphs[0].font.size = Pt(16)
season_frame.paragraphs[0].font.bold = True
season_frame.paragraphs[0].font.color.rgb = BRAND_BLACK
for p in season_frame.paragraphs[1:]:
    p.font.size = Pt(12)
    p.font.color.rgb = BRAND_DARK_GRAY

add_speaker_notes(
    slide4,
    "Why Kutch? Rann Utsav brings 5-7 Lakh visitors in just 4 months. Bhuj Railway Station and Gandhidham Junction are major transit hubs. Gujarat tourism is growing at 12% CAGR. We have the perfect test market with high tourist traffic, limited competition, and growing infrastructure. This is our beachhead market before pan-India expansion.",
)

print("Slides 1-4 created...")

# Save intermediate
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4.pptx"
)
print("Saved first 4 slides")
