#!/usr/bin/env python3
"""
BagDrop Presentation V5 - Part 2: Slides 15-20
Fixes: Better contrast on yellow backgrounds (BLACK text, not gray)
Swiss Style: Yellow/Black/White
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Load existing presentation
prs = Presentation(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V5.pptx"
)

# Colors
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 215, 0)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_GRAY = RGBColor(120, 120, 120)
GREEN = RGBColor(52, 168, 83)
RED = RGBColor(234, 67, 53)

HEADING_FONT = "Arial Black"
BODY_FONT = "Montserrat"


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_heading(
    slide, left, top, width, height, text, font_size, color, align=PP_ALIGN.LEFT
):
    return add_text(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size,
        color,
        bold=True,
        font=HEADING_FONT,
        align=align,
    )


def set_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================
# SLIDE 15: RISK MITIGATION - FIXED CONTRAST
# ============================================
s15 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s15, WHITE)
add_rect(s15, Inches(0), Inches(0), Inches(0.5), Inches(7.5), BLACK)
add_heading(
    s15, Inches(1), Inches(0.4), Inches(11), Inches(0.8), "RISK MITIGATION", 56, BLACK
)

risks = [
    ("SECURITY", "₹25K insurance per bag, CCTV monitoring at all partners"),
    ("THEFT/DAMAGE", "Instant claim processing, 48-hour resolution guarantee"),
    ("LOW ADOPTION", "Launch during Rann Utsav peak, hotel partnerships"),
    ("TECH FAILURE", "99.9% uptime SLA, offline OTP fallback system"),
    ("COMPETITION", "First-mover in Gujarat, unique delivery USP"),
]

for i, (risk, mitigation) in enumerate(risks):
    y = Inches(1.5 + i * 1.05)

    # Risk label - BLACK background with YELLOW text
    add_rect(s15, Inches(1), y, Inches(3.2), Inches(0.9), BLACK)
    add_text(
        s15,
        Inches(1.2),
        y + Inches(0.25),
        Inches(2.8),
        Inches(0.5),
        risk,
        16,
        YELLOW,
        bold=True,
    )

    # Mitigation - Light background with BLACK text
    add_rect(s15, Inches(4.2), y, Inches(8.6), Inches(0.9), RGBColor(245, 245, 245))
    add_text(
        s15,
        Inches(4.4),
        y + Inches(0.25),
        Inches(8.2),
        Inches(0.5),
        mitigation,
        14,
        BLACK,  # BLACK text, not DARK_GRAY
    )

add_notes(
    s15,
    "5 key risks addressed: Security (₹25K insurance), theft, adoption, tech, and competition. Each has a clear mitigation strategy.",
)

# ============================================
# SLIDE 16: EXIT STRATEGY - FIXED CONTRAST
# ============================================
s16 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s16, BLACK)
add_rect(s16, Inches(0), Inches(0), Inches(0.4), Inches(7.5), YELLOW)
add_heading(
    s16, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "EXIT STRATEGY", 56, YELLOW
)
add_text(
    s16,
    Inches(0.8),
    Inches(1.4),
    Inches(11),
    Inches(0.5),
    "Multiple paths to investor returns",
    24,
    WHITE,
)

exits = [
    ("STRATEGIC ACQUISITION", "MakeMyTrip, OYO, Yatra, ixigo", "₹50+ Cr"),
    ("SERIES A/B", "Venture capital growth round", "₹20-50 Cr"),
    ("FRANCHISE MODEL", "License to other cities/countries", "Ongoing Royalties"),
]

for i, (exit_type, desc, value) in enumerate(exits):
    y = Inches(2.3 + i * 1.5)

    # Exit type card - YELLOW background with BLACK text
    add_rect(s16, Inches(0.8), y, Inches(5), Inches(1.3), YELLOW)
    add_text(
        s16,
        Inches(1),
        y + Inches(0.15),
        Inches(4.6),
        Inches(0.5),
        exit_type,
        20,
        BLACK,
        bold=True,
    )
    add_text(
        s16,
        Inches(1),
        y + Inches(0.7),
        Inches(4.6),
        Inches(0.5),
        desc,
        14,
        DARK_GRAY,  # DARK_GRAY readable on YELLOW
    )

    # Value - on black side
    add_rect(s16, Inches(6), y, Inches(6.5), Inches(1.3), RGBColor(30, 30, 30))
    add_text(
        s16,
        Inches(6.2),
        y + Inches(0.35),
        Inches(6.1),
        Inches(0.6),
        value,
        28,
        YELLOW,
        bold=True,
    )

# Timeline
add_rect(s16, Inches(0.8), Inches(6.3), Inches(12), Inches(0.8), YELLOW)
add_text(
    s16,
    Inches(1),
    Inches(6.45),
    Inches(11.6),
    Inches(0.5),
    "TARGET TIMELINE: Year 4-5 for strategic exit or Series B",
    18,
    BLACK,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_notes(
    s16,
    "3 exit paths: Strategic acquisition (₹50Cr+), Series A/B funding, or franchise model. Target: Year 4-5.",
)

# ============================================
# SLIDE 17: ROADMAP / MILESTONES
# ============================================
s17 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s17, WHITE)
add_heading(
    s17, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), "ROADMAP", 56, BLACK
)
add_text(
    s17,
    Inches(0.8),
    Inches(1.1),
    Inches(11),
    Inches(0.5),
    "Key milestones",
    20,
    LIGHT_GRAY,
)

# Timeline bar
add_rect(s17, Inches(0.5), Inches(2.2), Inches(12.333), Inches(0.15), BLACK)

milestones = [
    ("Q1 2026", "LAUNCH", "Bhuj & Gandhidham\n20 partners"),
    ("Q2 2026", "SCALE", "50 partners\n5K users"),
    ("Q4 2026", "PROFIT", "Break-even\n100 partners"),
    ("2027", "EXPAND", "Pan-Gujarat\n500 partners"),
    ("2028", "NATIONAL", "Major cities\nFranchise launch"),
]

for i, (time, title, desc) in enumerate(milestones):
    x = Inches(0.3 + i * 2.6)
    color = YELLOW if i % 2 == 0 else BLACK
    txt_color = BLACK if i % 2 == 0 else YELLOW
    desc_color = DARK_GRAY if i % 2 == 0 else WHITE

    # Dot on timeline
    add_rect(s17, x + Inches(1.1), Inches(2.05), Inches(0.25), Inches(0.45), color)

    # Card
    add_rect(s17, x, Inches(2.8), Inches(2.4), Inches(3.5), color)
    add_text(
        s17,
        x,
        Inches(2.95),
        Inches(2.4),
        Inches(0.4),
        time,
        14,
        txt_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s17,
        x,
        Inches(3.4),
        Inches(2.4),
        Inches(0.5),
        title,
        22,
        txt_color,
        bold=True,
        font=HEADING_FONT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s17,
        x,
        Inches(4),
        Inches(2.4),
        Inches(2),
        desc,
        14,
        desc_color,
        align=PP_ALIGN.CENTER,
    )

# Bottom stat
add_rect(s17, Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.7), BLACK)
add_text(
    s17,
    Inches(1),
    Inches(6.65),
    Inches(11.333),
    Inches(0.5),
    "3-YEAR GOAL: ₹3.82 Cr Revenue | 15 Lakh Bookings | 60 Team Members",
    18,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_notes(
    s17,
    "5 milestones: Q1 launch, Q2 scale to 50 partners, Q4 break-even, 2027 pan-Gujarat, 2028 national expansion.",
)

# ============================================
# SLIDE 18: GOVERNMENT SUPPORT - FIXED CONTRAST
# ============================================
s18 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s18, BLACK)
add_rect(s18, Inches(0), Inches(0), Inches(13.333), Inches(0.3), YELLOW)
add_heading(
    s18,
    Inches(0.8),
    Inches(0.6),
    Inches(11),
    Inches(0.8),
    "GOVERNMENT SUPPORT",
    48,
    YELLOW,
)

# Schemes - using WHITE background cards with BLACK text for readability
schemes = [
    ("STARTUP INDIA", "Tax benefits, self-certification, fast-track patents"),
    ("PMMY (MUDRA)", "Up to ₹10 Lakh loan for micro enterprises"),
    ("iSTART RAJASTHAN", "Incubation, mentorship, seed funding"),
    ("TOURISM GRANTS", "Gujarat Tourism incentives for tourism startups"),
]

for i, (scheme, benefit) in enumerate(schemes):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.8 + row * 2.2)

    # Card with dark background (NOT yellow) for contrast
    add_rect(s18, x, y, Inches(6.2), Inches(1.9), RGBColor(30, 30, 30))
    add_rect(s18, x, y, Inches(6.2), Inches(0.12), YELLOW)  # Yellow top accent
    add_text(
        s18,
        x + Inches(0.2),
        y + Inches(0.3),
        Inches(5.8),
        Inches(0.5),
        scheme,
        20,
        YELLOW,
        bold=True,
    )
    add_text(
        s18,
        x + Inches(0.2),
        y + Inches(0.9),
        Inches(5.8),
        Inches(0.8),
        benefit,
        16,
        WHITE,
    )

# Bottom highlight - YELLOW with BLACK text (FIXED)
add_rect(s18, Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.8), YELLOW)
add_text(
    s18,
    Inches(1),
    Inches(6.45),
    Inches(11.333),
    Inches(0.5),
    "BagDrop qualifies for multiple government startup schemes",
    18,
    BLACK,  # BLACK text on YELLOW - highly visible!
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_notes(
    s18,
    "Government support: Startup India, MUDRA loan, iSTART, and tourism grants. Multiple schemes applicable.",
)

# ============================================
# SLIDE 19: CALL TO ACTION - FIXED CONTRAST
# ============================================
s19 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s19, YELLOW)  # YELLOW background
add_rect(s19, Inches(0), Inches(0), Inches(0.5), Inches(7.5), BLACK)

# All text is BLACK on YELLOW for maximum visibility
add_heading(
    s19,
    Inches(1),
    Inches(0.5),
    Inches(11),
    Inches(1),
    "LET'S BUILD THIS",
    64,
    BLACK,  # BLACK text on YELLOW
)
add_heading(
    s19,
    Inches(1),
    Inches(1.5),
    Inches(11),
    Inches(1),
    "TOGETHER",
    64,
    BLACK,  # BLACK text on YELLOW
)

# Value props with BLACK text
props = [
    "First-mover advantage in Gujarat's ₹500 Cr market",
    "Unique forgotten luggage delivery - NO competitor has this",
    "Break-even in 10 months, profitable by Year 1 Q4",
    "Clear path to ₹50+ Cr valuation by Year 5",
]

for i, prop in enumerate(props):
    y = Inches(2.8 + i * 0.7)
    add_rect(s19, Inches(1), y, Inches(0.15), Inches(0.4), BLACK)
    add_text(
        s19,
        Inches(1.3),
        y,
        Inches(11),
        Inches(0.5),
        prop,
        18,
        BLACK,  # BLACK text on YELLOW
        bold=True,
    )

# Ask box - BLACK background
add_rect(s19, Inches(1), Inches(5.4), Inches(5.5), Inches(1.5), BLACK)
add_text(
    s19,
    Inches(1.2),
    Inches(5.6),
    Inches(5.1),
    Inches(0.4),
    "SEED FUNDING",
    18,
    WHITE,
)
add_heading(
    s19, Inches(1.2), Inches(6), Inches(5.1), Inches(0.8), "₹36 LAKH", 48, YELLOW
)

# Contact box
add_rect(s19, Inches(7), Inches(5.4), Inches(5.5), Inches(1.5), BLACK)
add_text(
    s19,
    Inches(7.2),
    Inches(5.6),
    Inches(5.1),
    Inches(0.4),
    "CONTACT US",
    18,
    WHITE,
)
add_text(
    s19,
    Inches(7.2),
    Inches(6),
    Inches(5.1),
    Inches(0.8),
    "teamback@bagdrop.in",
    24,
    YELLOW,
    bold=True,
)
add_notes(
    s19,
    "Call to action: ₹36 Lakh seed funding. First-mover, unique USP, profitable by Year 1. Let's build this together!",
)

# ============================================
# SLIDE 20: THANK YOU / Q&A
# ============================================
s20 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s20, BLACK)
add_rect(s20, Inches(0), Inches(2.5), Inches(3), Inches(2.5), YELLOW)

add_heading(
    s20,
    Inches(2),
    Inches(2),
    Inches(10),
    Inches(1.5),
    "THANK YOU",
    96,
    WHITE,
)

add_text(
    s20,
    Inches(2),
    Inches(4),
    Inches(10),
    Inches(0.8),
    "Questions & Answers",
    32,
    YELLOW,
    bold=True,
)

add_text(
    s20,
    Inches(2),
    Inches(5.2),
    Inches(10),
    Inches(0.5),
    "TeamBack | Kutch, Gujarat | January 2026",
    20,
    LIGHT_GRAY,
)

# Contact info at bottom
add_rect(s20, Inches(0), Inches(6.5), Inches(13.333), Inches(1), YELLOW)
add_text(
    s20,
    Inches(0.5),
    Inches(6.7),
    Inches(6),
    Inches(0.5),
    "teamback@bagdrop.in",
    18,
    BLACK,
    bold=True,
)
add_text(
    s20,
    Inches(7),
    Inches(6.7),
    Inches(6),
    Inches(0.5),
    "www.bagdrop.in",
    18,
    BLACK,
    bold=True,
    align=PP_ALIGN.RIGHT,
)
add_notes(
    s20,
    "Thank you for your time. We're excited to answer any questions about BagDrop. Contact: teamback@bagdrop.in",
)

# Save
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V5.pptx"
)
print("Slides 15-20 added successfully!")
print("V5 Complete: 20 slides saved to BagDrop_Presentation_V5.pptx")
