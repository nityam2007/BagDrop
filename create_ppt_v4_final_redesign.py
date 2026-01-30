#!/usr/bin/env python3
"""
BagDrop Presentation V4 - FINAL REDESIGNED VERSION
Swiss Style: Bold Typography, Yellow Accent, Black/White Theme
Fonts: Arial Black (headings), Montserrat (body)
Complete 20-slide business plan presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Swiss Style Colors
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 215, 0)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_GRAY = RGBColor(120, 120, 120)

# Fonts
HEADING_FONT = "Arial Black"
BODY_FONT = "Montserrat"


# Helper functions
def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_heading(
    slide, left, top, width, height, text, font_size, color, align=PP_ALIGN.LEFT
):
    return add_text(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size,
        color,
        bold=True,
        font=HEADING_FONT,
        align=align,
    )


def set_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================
# SLIDE 1: TITLE
# ============================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, WHITE)
add_rect(s1, Inches(0), Inches(2), Inches(2), Inches(3), YELLOW)
add_heading(s1, Inches(1.5), Inches(2), Inches(10), Inches(1.5), "BAGDROP", 96, BLACK)
add_text(
    s1,
    Inches(1.5),
    Inches(3.5),
    Inches(10),
    Inches(0.8),
    "SMART LUGGAGE STORAGE & DELIVERY",
    32,
    DARK_GRAY,
    bold=True,
)
add_text(
    s1,
    Inches(1.5),
    Inches(4.5),
    Inches(10),
    Inches(0.5),
    "Secure. Contactless. Affordable.",
    24,
    YELLOW,
)
add_text(
    s1,
    Inches(1.5),
    Inches(6),
    Inches(10),
    Inches(0.5),
    "TeamBack | Kutch, Gujarat | January 2026",
    16,
    LIGHT_GRAY,
)
add_notes(
    s1,
    "Welcome judges! We're TeamBack from Kutch, Gujarat. BagDrop is India's smart luggage storage network. We solve a problem faced by 70% of travelers.",
)

# ============================================
# SLIDE 2: THE PROBLEM
# ============================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, WHITE)
add_rect(s2, Inches(10.5), Inches(0), Inches(2.833), Inches(1.8), YELLOW)
add_heading(
    s2, Inches(0.8), Inches(0.8), Inches(9), Inches(1), "THE PROBLEM", 64, BLACK
)
add_text(
    s2,
    Inches(0.8),
    Inches(1.8),
    Inches(9),
    Inches(0.6),
    "Travelers face a daily struggle with luggage",
    22,
    LIGHT_GRAY,
)

# Problem cards
add_rect(
    s2, Inches(0.5), Inches(2.8), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245)
)
add_text(
    s2,
    Inches(0.8),
    Inches(3.1),
    Inches(5.4),
    Inches(1),
    "EARLY CHECKOUT, LATE DEPARTURE\nHotels force checkout but flights are hours away",
    18,
    DARK_GRAY,
)

add_rect(
    s2, Inches(6.8), Inches(2.8), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245)
)
add_text(
    s2,
    Inches(7.1),
    Inches(3.1),
    Inches(5.4),
    Inches(1),
    "HEAVY BAGS, RUINED PLANS\nDragging luggage kills the travel experience",
    18,
    DARK_GRAY,
)

add_rect(
    s2, Inches(0.5), Inches(4.5), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245)
)
add_text(
    s2,
    Inches(0.8),
    Inches(4.8),
    Inches(5.4),
    Inches(1),
    "CLOAKROOMS? OVERPRICED\nRailway lockers cost ₹100+/day, often full",
    18,
    DARK_GRAY,
)

add_rect(
    s2, Inches(6.8), Inches(4.5), Inches(5.8), Inches(1.4), RGBColor(245, 245, 245)
)
add_text(
    s2,
    Inches(7.1),
    Inches(4.8),
    Inches(5.4),
    Inches(1),
    "FORGOT YOUR BAG?\nNo easy way to retrieve forgotten luggage",
    18,
    DARK_GRAY,
)

add_text(
    s2,
    Inches(0.8),
    Inches(6.3),
    Inches(11),
    Inches(0.5),
    "70% of travelers have faced luggage hassles during their trip",
    20,
    YELLOW,
    bold=True,
)
add_notes(
    s2,
    "70% of travelers face luggage problems. Early checkout with late departure means hours of frustration. Railway lockers are expensive at ₹100+/day and often unavailable.",
)

# ============================================
# SLIDE 3: THE SOLUTION
# ============================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, BLACK)
add_rect(s3, Inches(0), Inches(0), Inches(0.3), Inches(7.5), YELLOW)
add_heading(
    s3, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "THE SOLUTION", 64, YELLOW
)
add_text(
    s3,
    Inches(0.8),
    Inches(1.6),
    Inches(11),
    Inches(1),
    "Store your bags anywhere. Pick them up anytime.",
    32,
    WHITE,
    bold=True,
)

add_rect(s3, Inches(0.8), Inches(3), Inches(0.15), Inches(0.8), YELLOW)
add_text(
    s3,
    Inches(1.2),
    Inches(3),
    Inches(10),
    Inches(0.9),
    "NEIGHBORHOOD NETWORK\nWe partner with verified local shops, cafes & hotels as storage points",
    20,
    WHITE,
)

add_rect(s3, Inches(0.8), Inches(4.2), Inches(0.15), Inches(0.8), YELLOW)
add_text(
    s3,
    Inches(1.2),
    Inches(4.2),
    Inches(10),
    Inches(0.9),
    "INSTANT ACCESS\nBook via app, drop your bag, get a secure OTP or NFC token",
    20,
    WHITE,
)

add_rect(s3, Inches(0.8), Inches(5.4), Inches(0.15), Inches(0.8), YELLOW)
add_text(
    s3,
    Inches(1.2),
    Inches(5.4),
    Inches(10),
    Inches(0.9),
    "FORGOTTEN LUGGAGE? DELIVERED\nLeft something behind? We'll courier it anywhere in India",
    20,
    WHITE,
)
add_notes(
    s3,
    "BagDrop partners with local shops for secure storage. Contactless OTP/NFC access. Unique delivery service for forgotten luggage - NO competitor offers this.",
)

# ============================================
# SLIDE 4: WHY GUJARAT
# ============================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, YELLOW)
add_heading(
    s4, Inches(0.8), Inches(0.5), Inches(11), Inches(1), "WHY GUJARAT?", 64, BLACK
)

# Kutch advantages
advantages = [
    ("TOURIST HUB", "Rann Utsav: 5L+ visitors annually"),
    ("TRANSIT", "Bhuj Railway, Gandhidham Junction"),
    ("GROWTH", "12% CAGR tourism growth"),
    ("STRATEGIC", "Gateway to Rann, Mundra Port"),
    ("CULTURE", "Dholavira, handicrafts, year-round"),
]

for i, (title, desc) in enumerate(advantages):
    y = Inches(1.8 + i * 0.9)
    add_text(s4, Inches(0.8), y, Inches(4), Inches(0.4), title, 20, BLACK, bold=True)
    add_text(
        s4, Inches(0.8), y + Inches(0.4), Inches(6), Inches(0.4), desc, 16, DARK_GRAY
    )

# Stats
add_rect(s4, Inches(8), Inches(2), Inches(4.5), Inches(2), BLACK)
add_text(
    s4,
    Inches(8.2),
    Inches(2.3),
    Inches(4.1),
    Inches(0.6),
    "5.5 CRORE",
    36,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_text(
    s4,
    Inches(8.2),
    Inches(3),
    Inches(4.1),
    Inches(0.5),
    "Annual Gujarat tourists",
    14,
    WHITE,
    align=PP_ALIGN.CENTER,
)

add_rect(s4, Inches(8), Inches(4.5), Inches(4.5), Inches(1.5), BLACK)
add_text(
    s4,
    Inches(8.2),
    Inches(4.7),
    Inches(4.1),
    Inches(0.5),
    "12% CAGR",
    28,
    YELLOW,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_text(
    s4,
    Inches(8.2),
    Inches(5.3),
    Inches(4.1),
    Inches(0.4),
    "Tourism growth rate",
    14,
    WHITE,
    align=PP_ALIGN.CENTER,
)

add_notes(
    s4,
    "Kutch, Gujarat is perfect for launch. Rann Utsav brings 5L+ visitors. Major railway stations. 12% tourism growth. Strategic location near Mundra Port and Dholavira.",
)

# ============================================
# SLIDE 5: TARGET MARKET
# ============================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, WHITE)
add_rect(s5, Inches(0), Inches(0), Inches(0.5), Inches(7.5), YELLOW)
add_heading(
    s5, Inches(1), Inches(0.5), Inches(11), Inches(1), "TARGET MARKET", 64, BLACK
)

audience = """-> TOURISTS & BACKPACKERS
   Rann Utsav visitors, sightseers

-> TRAIN & BUS TRAVELERS
   Inter-city commuters, layovers

-> STUDENTS & DAILY COMMUTERS
   Budget-conscious, UPI users

-> BUSINESS TRAVELERS
   Premium service, NFC preferred

-> FORGOTTEN LUGGAGE USERS
   High-value delivery service"""

add_text(s5, Inches(1), Inches(2), Inches(11), Inches(5), audience, 28, DARK_GRAY)
add_notes(
    s5,
    "We serve 5 segments: Tourists visiting Rann Utsav, train/bus travelers with layovers, budget students, business travelers needing premium service, and our unique segment - forgotten luggage customers.",
)

# ============================================
# SLIDE 6: MARKET SIZE
# ============================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6, BLACK)
add_heading(
    s6,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.8),
    "MARKET OPPORTUNITY",
    56,
    YELLOW,
)

# Stats boxes
stats = [
    ("TAM", "₹2,000+ Cr", "Total Addressable Market"),
    ("SAM", "₹500 Cr", "Serviceable (Gujarat + 10 cities)"),
    ("SOM", "₹36 Lakh", "Obtainable (Year 1 target)"),
]

for i, (label, value, desc) in enumerate(stats):
    x = Inches(0.5 + i * 4.2)
    add_rect(
        s6,
        x,
        Inches(1.8),
        Inches(4),
        Inches(2.2),
        YELLOW if i == 2 else RGBColor(40, 40, 40),
    )
    color = BLACK if i == 2 else YELLOW
    add_text(
        s6,
        x,
        Inches(2),
        Inches(4),
        Inches(0.8),
        value,
        36,
        color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s6,
        x,
        Inches(2.9),
        Inches(4),
        Inches(0.5),
        label,
        18,
        color,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s6,
        x,
        Inches(3.4),
        Inches(4),
        Inches(0.5),
        desc,
        12,
        LIGHT_GRAY if i < 2 else DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )

# Market data
add_text(
    s6,
    Inches(0.8),
    Inches(4.5),
    Inches(5),
    Inches(2),
    "GUJARAT:\n5.5 Cr domestic tourists\n5.2 Lakh foreign tourists\n12% CAGR growth",
    18,
    WHITE,
)

add_text(
    s6,
    Inches(7),
    Inches(4.5),
    Inches(5.5),
    Inches(2),
    "INDIA:\n2,400 Cr railway passengers\n15 Cr air passengers\n₹15 Lakh Cr tourism GDP",
    18,
    WHITE,
)

add_notes(
    s6,
    "Massive market opportunity. TAM is ₹2,000+ Crore. SAM of ₹500 Crore in Gujarat + major cities. Our Year 1 target is ₹36 Lakh, growing to ₹1.27 Crore by Year 2.",
)

# ============================================
# SLIDE 7: BUSINESS MODEL
# ============================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7, WHITE)
add_heading(
    s7, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8), "BUSINESS MODEL", 56, BLACK
)
add_text(
    s7,
    Inches(0.8),
    Inches(1.1),
    Inches(11),
    Inches(0.5),
    "Multiple revenue streams",
    20,
    LIGHT_GRAY,
)

# Revenue streams table
revenue = [
    ("STORAGE FEES", "30% commission", "₹50-200 per booking"),
    ("DELIVERY FEES", "20% commission", "₹300-1,500 forgotten bags"),
    ("NFC TOKENS", "100% margin", "₹199 one-time"),
    ("KIOSK ADS", "100% margin", "₹5,000/month per kiosk"),
    ("PARTNER SUBS", "100% margin", "₹500/month premium"),
]

# Headers
headers = ["REVENUE STREAM", "MARGIN", "PRICING"]
for col, header in enumerate(headers):
    x = Inches(0.5 + col * 4)
    add_rect(s7, x, Inches(1.8), Inches(4), Inches(0.4), BLACK)
    add_text(
        s7,
        x,
        Inches(1.85),
        Inches(4),
        Inches(0.35),
        header,
        13,
        YELLOW,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

# Data
for row_idx, (stream, margin, price) in enumerate(revenue):
    y = Inches(2.2 + row_idx * 0.6)
    data = [stream, margin, price]
    for col, val in enumerate(data):
        x = Inches(0.5 + col * 4)
        bg_color = RGBColor(245, 245, 245) if row_idx % 2 == 0 else WHITE
        add_rect(s7, x, y, Inches(4), Inches(0.55), bg_color)
        txt_color = (
            BLACK if col == 0 else (RGBColor(52, 168, 83) if col == 1 else DARK_GRAY)
        )
        add_text(
            s7,
            x + Inches(0.2),
            y + Inches(0.15),
            Inches(3.6),
            Inches(0.3),
            val,
            11,
            txt_color,
        )

# Pricing highlight
add_rect(s7, Inches(0.5), Inches(5.5), Inches(6), Inches(1.5), YELLOW)
add_text(
    s7,
    Inches(0.7),
    Inches(5.7),
    Inches(5.6),
    Inches(0.6),
    "AFFORDABLE PRICING",
    20,
    BLACK,
    bold=True,
)
add_text(
    s7,
    Inches(0.7),
    Inches(6.3),
    Inches(5.6),
    Inches(0.5),
    "Starting at ₹50 for 6 hours",
    18,
    DARK_GRAY,
)
add_text(
    s7,
    Inches(0.7),
    Inches(6.7),
    Inches(5.6),
    Inches(0.3),
    "Cheaper than railway cloakrooms",
    12,
    LIGHT_GRAY,
)

add_notes(
    s7,
    "5 revenue streams. Primary is 30% commission on storage. 20% on delivery fees. NFC tokens, kiosk ads, and partner subscriptions add pure margin. Starting price ₹50 for 6 hours beats railway lockers.",
)

print("Slides 1-7 complete")

# Save progress
prs.save(
    "/home/nityam/Downloads/code/STUPID PROJECTs/BagDrop/BagDrop_Presentation_V4_Final.pptx"
)
print("Saved slides 1-7")
